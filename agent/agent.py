import os
import json
import time
import asyncio
import logging
import threading
from datetime import datetime
from typing import Dict, Any, List, Optional
import boto3
from botocore.exceptions import ClientError
from bedrock_agentcore import BedrockAgentCoreApp, PingStatus
from strands import Agent, tool
from strands_tools import http_request

# Platform Agents
from infrastructure_agent import (
    set_agents_config_reference as set_infra_config_reference,
    get_system_status,
    get_recent_errors,
    get_trace_summary,
    check_model_availability,
    check_agent_tools,
    INFRASTRUCTURE_AGENT_PROMPT,
)
from security_agent import (
    set_agents_config_reference as set_security_config_reference,
    audit_iam_role,
    check_bedrock_guardrails,
    check_secrets_compliance,
    scan_text_for_pii,
    SECURITY_AGENT_PROMPT,
)
from observability_agent import (
    set_agents_config_reference as set_observability_config_reference,
    get_runtime_metrics,
    get_session_details,
    get_user_metrics,
    OBSERVABILITY_AGENT_PROMPT,
)
from report_tool import generate_report

# AgentCore Memory imports
from bedrock_agentcore.memory import MemoryClient
from bedrock_agentcore.memory.integrations.strands.config import (
    AgentCoreMemoryConfig,
    RetrievalConfig
)
from bedrock_agentcore.memory.integrations.strands.session_manager import (
    AgentCoreMemorySessionManager
)

# =============================================================================
# LOGGING CONFIGURATION
# =============================================================================
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# =============================================================================
# INITIALIZE AGENTCORE APP
# =============================================================================
app = BedrockAgentCoreApp()

# =============================================================================
# ENVIRONMENT VARIABLES
# =============================================================================
USER_ID = os.environ.get('USER_ID', '202')
PROJECT_ID = os.environ.get('PROJECT_ID', '202')
DYNAMODB_TABLE_NAME = os.environ.get('DYNAMODB_TABLE_NAME', 'AgentConfigs')
AWS_REGION = os.environ.get('AWS_REGION', 'eu-central-1')
COORDINATOR_KEY = os.environ.get('COORDINATOR_KEY', 'central_coordinator')

# AgentCore Memory configuration
# Memory ID must match pattern: [a-zA-Z][a-zA-Z0-9-_]{0,99}-[a-zA-Z0-9]{10}
# Example: ArchitectureAnalyzerMemory-QXesnp3yGd
# MUST be set via AGENTCORE_MEMORY_ID environment variable (set by CDK deployment)
MEMORY_ID = os.environ.get('AGENTCORE_MEMORY_ID', '')
MEMORY_REGION = os.environ.get('MEMORY_REGION', 'eu-central-1')
ENABLE_SHORT_TERM_MEMORY = os.environ.get('ENABLE_SHORT_TERM_MEMORY', 'true').lower() == 'true'
ENABLE_LONG_TERM_MEMORY = os.environ.get('ENABLE_LONG_TERM_MEMORY', 'true').lower() == 'true'

# S3 Vectors configuration
S3_DOCS_BUCKET = os.environ.get('S3_DOCS_BUCKET', 'qubitz-customer-prod-v2')
S3_VECTOR_BUCKET = os.environ.get('S3_VECTOR_BUCKET', 'qubitz-vectors-prod')

# Amplify buckets for agent-isolated uploads
AMPLIFY_BUCKET_DEV = os.environ.get('AMPLIFY_BUCKET_DEV', 'amplify-d2modhwpvv0amc-de-qubitzamplifybucketbucke-ndovyciox3mg')
AMPLIFY_BUCKET_MAIN = os.environ.get('AMPLIFY_BUCKET_MAIN', 'amplify-d2modhwpvv0amc-ma-qubitzamplifybucketbucke-aaekcylhcq1i')

EMBEDDING_MODEL_ID = os.environ.get('EMBEDDING_MODEL_ID', 'amazon.titan-embed-text-v2:0')
EMBEDDING_DIMENSION = int(os.environ.get('EMBEDDING_DIMENSION', '1024'))
EMBEDDING_REGION = os.environ.get('EMBEDDING_REGION', 'eu-central-1')

# Gateway configuration for MCP-based tools (web search via Tavily/SerpAPI, custom MCP servers)
GATEWAY_URL = os.environ.get('GATEWAY_URL')

# Credentials table for per-user connector credentials (used by gateway Lambdas)
CREDENTIALS_TABLE_NAME = os.environ.get('CREDENTIALS_TABLE_NAME', 'ToolCredentials')

# Connector configuration (Identity-based OAuth — direct MCP, no gateway)
CONNECTOR_TARGETS_JSON = os.environ.get('CONNECTOR_TARGETS', '[]')
try:
    CONNECTOR_TARGETS = json.loads(CONNECTOR_TARGETS_JSON) if CONNECTOR_TARGETS_JSON else []
except (json.JSONDecodeError, TypeError):
    CONNECTOR_TARGETS = []
# Build lookup: connector_name → config
_CONNECTOR_TARGET_MAP = {t.get('name', ''): t for t in CONNECTOR_TARGETS if t.get('name')}


def _fix_model_region(model_id: str) -> str:
    """Ensure model ID uses the correct region prefix for the current AWS_REGION."""
    if not model_id:
        return model_id
    region_prefix = 'eu' if AWS_REGION.startswith('eu') else 'us'
    # Replace wrong region prefix (e.g. us. -> eu. when in eu-central-1)
    if model_id.startswith('us.') and region_prefix == 'eu':
        model_id = 'eu.' + model_id[3:]
    elif model_id.startswith('eu.') and region_prefix == 'us':
        model_id = 'us.' + model_id[3:]
    return model_id


# =============================================================================
# AWS CLIENTS
# =============================================================================
dynamodb = boto3.resource('dynamodb', region_name=AWS_REGION)
table = dynamodb.Table(DYNAMODB_TABLE_NAME)
s3_client = boto3.client('s3', region_name=AWS_REGION)
s3vectors_client = boto3.client('s3vectors', region_name=EMBEDDING_REGION)
bedrock_runtime = boto3.client('bedrock-runtime', region_name=EMBEDDING_REGION)


# =============================================================================
# S3 LOGGER
# =============================================================================
class S3Logger:
    """Logger that appends logs to S3 at {bucket}/{user_id}/{project_id}/logs.txt"""
    
    def __init__(self, bucket: str, user_id: str, project_id: str):
        self.bucket = bucket
        self.user_id = user_id
        self.project_id = project_id
        self.log_key = f"{user_id}/{project_id}/logs.txt"
        self._buffer = []
    
    def log(self, level: str, message: str):
        """Add a log entry to buffer"""
        timestamp = datetime.utcnow().isoformat()
        entry = f"[{timestamp}] [{level}] {message}"
        self._buffer.append(entry)
        getattr(logger, level.lower(), logger.info)(message)
    
    def info(self, message: str):
        self.log("INFO", message)
    
    def error(self, message: str):
        self.log("ERROR", message)
    
    def warning(self, message: str):
        self.log("WARNING", message)
    
    def flush(self):
        """Flush buffer to S3"""
        if not self._buffer:
            return
        
        try:
            existing_content = ""
            try:
                response = s3_client.get_object(Bucket=self.bucket, Key=self.log_key)
                existing_content = response['Body'].read().decode('utf-8')
            except ClientError as e:
                if e.response['Error']['Code'] != 'NoSuchKey':
                    raise
            
            new_content = existing_content + "\n".join(self._buffer) + "\n"
            
            s3_client.put_object(
                Bucket=self.bucket,
                Key=self.log_key,
                Body=new_content.encode('utf-8'),
                ContentType='text/plain'
            )
            
            logger.info(f"Flushed {len(self._buffer)} log entries to s3://{self.bucket}/{self.log_key}")
            self._buffer = []
        except Exception as e:
            logger.error(f"Failed to flush logs to S3: {str(e)}")


# Global S3 logger instance
s3_logger: Optional[S3Logger] = None


def get_s3_logger() -> S3Logger:
    """Get or create S3 logger instance"""
    global s3_logger
    if s3_logger is None:
        s3_logger = S3Logger(S3_DOCS_BUCKET, USER_ID, PROJECT_ID)
    return s3_logger


# =============================================================================
# AGENTCORE MEMORY MANAGER (FIXED)
# =============================================================================
class AgentCoreMemoryManager:
    """
    Manages AgentCore Memory for agents with both short-term and long-term capabilities.
    
    CRITICAL FIXES:
    - Session ID and Actor ID must be provided per request for proper isolation
    - Proper namespace patterns with template variables
    - Session manager created per request with correct IDs
    """
    
    def __init__(
        self,
        memory_id: Optional[str] = None,
        region: str = 'eu-central-1',
        enable_short_term: bool = True,
        enable_long_term: bool = True
    ):
        self.memory_id = memory_id
        self.region = region
        self.enable_short_term = enable_short_term
        self.enable_long_term = enable_long_term
        self.memory_client = MemoryClient(region_name=region) if memory_id else None
        
        if memory_id:
            logger.info(f"AgentCore Memory initialized with ID: {memory_id}")
            logger.info(f"Short-term memory: {enable_short_term}, Long-term memory: {enable_long_term}")
        else:
            logger.warning("Memory ID not configured. Memory features disabled.")
    
    def create_session_manager(
        self,
        actor_id: str,
        session_id: str,
        retrieval_config: Optional[Dict[str, RetrievalConfig]] = None
    ) -> Optional[AgentCoreMemorySessionManager]:
        """
        Create a session manager for an agent with memory capabilities.
        
        CRITICAL NOTES:
        - actor_id: Unique user identifier (for user isolation) - MUST come from payload
        - session_id: Conversation session ID - MUST persist across requests in the same conversation
        - retrieval_config: Namespace configuration for long-term memory retrieval
        
        Args:
            actor_id: Unique identifier for the user (e.g., from JWT token or user_id in payload)
            session_id: Unique identifier for conversation session (from payload or generated once)
            retrieval_config: Optional configuration for retrieving from different namespaces
            
        Returns:
            AgentCoreMemorySessionManager or None if memory is disabled
        """
        if not self.memory_id:
            logger.warning("Memory ID not configured. Memory features disabled.")
            return None
        
        # Build retrieval config for long-term memory namespaces
        # CRITICAL: Use template variables {actorId} and {sessionId} - NOT hard-coded values
        if self.enable_long_term and retrieval_config is None:
            retrieval_config = {
                # User preferences namespace - retrieves user-specific preferences
                "/preferences/{actorId}": RetrievalConfig(
                    top_k=5,
                    relevance_score=0.7
                ),
                # Facts namespace - retrieves user-specific facts
                "/facts/{actorId}": RetrievalConfig(
                    top_k=10,
                    relevance_score=0.5
                ),
                # Session summaries namespace - retrieves session-specific summaries
                "/summaries/{actorId}/{sessionId}": RetrievalConfig(
                    top_k=3,
                    relevance_score=0.6
                )
            }
        
        # Create memory config with proper IDs
        config = AgentCoreMemoryConfig(
            memory_id=self.memory_id,
            session_id=session_id,
            actor_id=actor_id,
            retrieval_config=retrieval_config if self.enable_long_term else None
        )
        
        # Create and return session manager
        session_manager = AgentCoreMemorySessionManager(
            agentcore_memory_config=config,
            region_name=self.region
        )
        
        logger.info(f"[OK] Session manager created: actor={actor_id}, session={session_id}")
        logger.info(f"  Short-term: {self.enable_short_term}, Long-term: {self.enable_long_term}")
        if retrieval_config:
            logger.info(f"  Retrieval namespaces: {list(retrieval_config.keys())}")
        
        return session_manager
    
    def get_memory_info(self) -> Dict[str, Any]:
        """Get information about the memory configuration"""
        return {
            "memory_id": self.memory_id,
            "region": self.region,
            "short_term_enabled": self.enable_short_term,
            "long_term_enabled": self.enable_long_term,
            "memory_configured": self.memory_id is not None
        }


# Global memory manager instance
memory_manager: Optional[AgentCoreMemoryManager] = None


def get_memory_manager() -> AgentCoreMemoryManager:
    """Get or create memory manager instance"""
    global memory_manager
    if memory_manager is None:
        memory_manager = AgentCoreMemoryManager(
            memory_id=MEMORY_ID,
            region=MEMORY_REGION,
            enable_short_term=ENABLE_SHORT_TERM_MEMORY,
            enable_long_term=ENABLE_LONG_TERM_MEMORY
        )
    return memory_manager


# =============================================================================
# S3 VECTORS RAG
# =============================================================================
class S3VectorsRAG:
    """
    RAG implementation using S3 Vectors for document retrieval.
    Documents are stored at: qubitz-customer-prod-v2/{user_id}/{project_id}/
    
    Loads configuration from DynamoDB including:
    - data_sources: List of configured connectors (S3, web_crawler, confluence, sharepoint, salesforce)
    - embedding_model: Model to use for embeddings
    """
    
    def __init__(self, user_id: str, project_id: str):
        self.user_id = user_id
        self.project_id = project_id
        self.docs_prefix = f"{user_id}/{project_id}/"
        sanitized_user = user_id.replace('_', '-').lower()
        sanitized_project = project_id.replace('_', '-').lower()
        self.index_name = f"{sanitized_user}-{sanitized_project}-docs"
        self._index_initialized = False
        
        # Load KB config from DynamoDB
        self.kb_config = self._load_kb_config_from_dynamodb()
        
        logger.info(f"RAG initialized with index name: {self.index_name}")
        if self.kb_config:
            data_sources = self.kb_config.get('data_sources', [])
            logger.info(f"KB config loaded: {len(data_sources)} data source(s) configured")
    
    def _load_kb_config_from_dynamodb(self) -> Dict[str, Any]:
        """Load knowledge base configuration from DynamoDB"""
        try:
            response = table.get_item(
                Key={
                    'user_id': self.user_id,
                    'project_id': self.project_id
                }
            )
            
            if 'Item' not in response:
                logger.warning(f"No DynamoDB item found for {self.user_id}/{self.project_id}")
                return {}
            
            item = response['Item']
            
            # Extract kb_config from response
            if 'response' in item and 'kb_config' in item['response']:
                kb_config_raw = item['response']['kb_config']
                
                # Parse data_sources if it's a JSON string
                data_sources_raw = kb_config_raw.get('data_sources', '[]')
                if isinstance(data_sources_raw, str):
                    try:
                        data_sources = json.loads(data_sources_raw)
                    except json.JSONDecodeError:
                        data_sources = []
                else:
                    data_sources = data_sources_raw
                
                kb_config = {
                    'data_sources': data_sources,
                    'embedding_model': kb_config_raw.get('embedding_model', 'amazon.titan-embed-text-v2:0'),
                    'start_sync': kb_config_raw.get('start_sync', True)
                }
                
                logger.info(f"Loaded KB config from DynamoDB: {len(data_sources)} data sources")
                return kb_config
            else:
                logger.info(f"No kb_config found in DynamoDB for {self.user_id}/{self.project_id}")
                return {}
                
        except Exception as e:
            logger.error(f"Error loading KB config from DynamoDB: {str(e)}")
            return {}
    
    def get_configured_data_sources(self) -> List[Dict[str, Any]]:
        """Get list of configured data sources from DynamoDB"""
        return self.kb_config.get('data_sources', [])
    
    def get_embedding_model(self) -> str:
        """Get configured embedding model"""
        return self.kb_config.get('embedding_model', EMBEDDING_MODEL_ID)
    
    def list_agent_uploads(self, agent_name: str, tool_name: str = 'rag', include_amplify_buckets: bool = True) -> List[Dict[str, Any]]:
        """
        List agent-isolated uploads from Amplify buckets.
        
        Path structure: user-uploads/{user_id}/{project_id}/{agent_name}/{tool_name}/
        
        Args:
            agent_name: Name of the agent (e.g., 'space_utilization_analyst')
            tool_name: Tool name (default: 'rag')
            include_amplify_buckets: Whether to check Amplify buckets (default: True)
        
        Returns:
            List of documents with metadata including bucket and key
        """
        documents = []
        upload_prefix = f"user-uploads/{self.user_id}/{self.project_id}/{agent_name}/{tool_name}/"
        
        buckets_to_check = []
        if include_amplify_buckets:
            buckets_to_check = [AMPLIFY_BUCKET_DEV, AMPLIFY_BUCKET_MAIN]
        
        # Also check main docs bucket
        buckets_to_check.append(S3_DOCS_BUCKET)
        
        for bucket in buckets_to_check:
            try:
                response = s3_client.list_objects_v2(
                    Bucket=bucket,
                    Prefix=upload_prefix
                )
                
                for obj in response.get('Contents', []):
                    key = obj['Key']
                    if key != upload_prefix:  # Skip the prefix itself
                        documents.append({
                            'bucket': bucket,
                            'key': key,
                            'name': key.replace(upload_prefix, ''),
                            'size': obj['Size'],
                            'last_modified': obj['LastModified'].isoformat(),
                            'agent_name': agent_name,
                            'tool_name': tool_name
                        })
                
                logger.info(f"Found {len(response.get('Contents', []))} files in {bucket}/{upload_prefix}")
                
            except Exception as e:
                logger.debug(f"Could not list from {bucket}/{upload_prefix}: {str(e)}")
                continue
        
        logger.info(f"Total agent uploads found: {len(documents)} for {agent_name}/{tool_name}")
        return documents
    
    def get_agent_upload_content(self, bucket: str, doc_key: str) -> str:
        """
        Retrieve document content from agent-isolated upload.
        
        Args:
            bucket: S3 bucket name (Amplify or main docs bucket)
            doc_key: Full S3 key path
        
        Returns:
            Document content as string
        """
        try:
            response = s3_client.get_object(
                Bucket=bucket,
                Key=doc_key
            )
            
            file_bytes = response['Body'].read()
            doc_lower = doc_key.lower()
            
            # Handle different file types
            if doc_lower.endswith('.pdf'):
                return self._extract_pdf_text(file_bytes, doc_key)
            elif doc_lower.endswith('.docx'):
                return self._extract_docx_text(file_bytes, doc_key)
            elif doc_lower.endswith('.csv'):
                return self._extract_csv_text(file_bytes, doc_key)
            elif doc_lower.endswith(('.xlsx', '.xls')):
                return self._extract_excel_text(file_bytes, doc_key)
            elif doc_lower.endswith(('.pptx', '.ppt')):
                return self._extract_ppt_text(file_bytes, doc_key)
            else:
                # Try to decode as UTF-8 text
                return file_bytes.decode('utf-8')
                
        except UnicodeDecodeError as e:
            logger.warning(f"Cannot decode {doc_key} as text (binary file?): {str(e)}")
            return ""
        except Exception as e:
            logger.error(f"Error reading agent upload {bucket}/{doc_key}: {str(e)}")
            return ""
    
    def index_agent_uploads(self, agent_name: str, tool_name: str = 'rag') -> Dict[str, int]:
        """
        Index all agent-isolated uploads into the vector store.
        
        Args:
            agent_name: Name of the agent
            tool_name: Tool name (default: 'rag')
        
        Returns:
            Dict mapping document names to number of chunks indexed
        """
        documents = self.list_agent_uploads(agent_name, tool_name)
        results = {}
        
        for doc in documents:
            bucket = doc['bucket']
            doc_key = doc['key']
            doc_name = doc['name']
            
            # Skip non-indexable files
            skip_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.ico', '.svg', 
                              '.mp4', '.avi', '.mov', '.mp3', '.wav', '.zip', '.tar', '.gz']
            if any(doc_name.lower().endswith(ext) for ext in skip_extensions):
                logger.debug(f"Skipping non-indexable file: {doc_name}")
                results[doc_name] = 0
                continue
            
            try:
                # Get content
                content = self.get_agent_upload_content(bucket, doc_key)
                if not content:
                    results[doc_name] = 0
                    continue
                
                # Chunk content
                chunks = []
                chunk_size = 1000
                overlap = 200
                for i in range(0, len(content), chunk_size - overlap):
                    chunk = content[i:i + chunk_size]
                    if chunk.strip():
                        chunks.append(chunk)
                
                # Create embeddings and store vectors
                vectors = []
                for idx, chunk in enumerate(chunks):
                    try:
                        embedding = self._get_embedding(chunk)
                        vector_key = f"{doc_key}_chunk_{idx}"
                        vectors.append({
                            'key': vector_key,
                            'data': {'float32': embedding},
                            'metadata': {
                                'doc_key': doc_key,
                                'doc_name': doc_name,
                                'bucket': bucket,
                                'agent_name': agent_name,
                                'tool_name': tool_name,
                                'chunk_index': idx,
                                'content': chunk[:500]
                            }
                        })
                    except Exception as e:
                        logger.error(f"Error creating embedding for {doc_name} chunk {idx}: {e}")
                
                if vectors:
                    self._ensure_vector_index()
                    batch_size = 100
                    for i in range(0, len(vectors), batch_size):
                        batch = vectors[i:i + batch_size]
                        s3vectors_client.put_vectors(
                            vectorBucketName=S3_VECTOR_BUCKET,
                            indexName=self.index_name,
                            vectors=batch
                        )
                    logger.info(f"Indexed {len(vectors)} chunks from {doc_name} (agent: {agent_name})")
                    results[doc_name] = len(vectors)
                else:
                    results[doc_name] = 0
                    
            except Exception as e:
                logger.error(f"Error indexing agent upload {doc_name}: {e}")
                results[doc_name] = 0
        
        return results
    
    def _get_embedding(self, text: str, purpose: str = "GENERIC_INDEX") -> List[float]:
        """Generate embedding for text using Amazon Titan Embed Text v2"""
        try:
            request_body = {
                "inputText": text,
                "dimensions": EMBEDDING_DIMENSION,
                "normalize": True
            }
            
            response = bedrock_runtime.invoke_model(
                modelId=EMBEDDING_MODEL_ID,
                body=json.dumps(request_body),
                contentType="application/json",
            )
            
            result = json.loads(response['body'].read())
            return result['embedding']
        except Exception as e:
            logger.error(f"Error generating embedding: {str(e)}")
            raise
    
    def _ensure_vector_bucket(self) -> bool:
        """Ensure vector bucket exists, create if not"""
        try:
            s3vectors_client.get_vector_bucket(vectorBucketName=S3_VECTOR_BUCKET)
            logger.info(f"Vector bucket {S3_VECTOR_BUCKET} exists")
            return True
        except ClientError as e:
            if e.response['Error']['Code'] in ['NotFoundException', 'ResourceNotFoundException']:
                try:
                    s3vectors_client.create_vector_bucket(vectorBucketName=S3_VECTOR_BUCKET)
                    logger.info(f"Created vector bucket: {S3_VECTOR_BUCKET}")
                    return True
                except Exception as create_err:
                    logger.error(f"Error creating vector bucket: {str(create_err)}")
                    return False
            else:
                logger.error(f"Error checking vector bucket: {str(e)}")
                return False
    
    def _ensure_vector_index(self) -> bool:
        """Ensure vector index exists, create if not"""
        if self._index_initialized:
            return True
        
        if not self._ensure_vector_bucket():
            return False
        
        try:
            s3vectors_client.get_index(
                vectorBucketName=S3_VECTOR_BUCKET,
                indexName=self.index_name
            )
            self._index_initialized = True
            logger.info(f"Vector index {self.index_name} exists")
            return True
        except ClientError as e:
            if e.response['Error']['Code'] in ['NotFoundException', 'ResourceNotFoundException']:
                try:
                    s3vectors_client.create_index(
                        vectorBucketName=S3_VECTOR_BUCKET,
                        indexName=self.index_name,
                        dataType='float32',
                        dimension=EMBEDDING_DIMENSION,
                        distanceMetric='cosine'
                    )
                    self._index_initialized = True
                    logger.info(f"Created vector index: {self.index_name}")
                    return True
                except Exception as create_err:
                    logger.error(f"Error creating index: {str(create_err)}")
                    return False
            else:
                logger.error(f"Error checking index: {str(e)}")
                return False
    
    def list_documents(self) -> List[Dict[str, Any]]:
        """List all documents in the user's project docs folder"""
        try:
            response = s3_client.list_objects_v2(
                Bucket=S3_DOCS_BUCKET,
                Prefix=self.docs_prefix
            )
            
            documents = []
            for obj in response.get('Contents', []):
                key = obj['Key']
                if key != self.docs_prefix:
                    documents.append({
                        'key': key,
                        'name': key.replace(self.docs_prefix, ''),
                        'size': obj['Size'],
                        'last_modified': obj['LastModified'].isoformat()
                    })
            
            logger.info(f"Found {len(documents)} documents in {self.docs_prefix}")
            return documents
        except Exception as e:
            logger.error(f"Error listing documents: {str(e)}")
            return []
    
    def get_document_content(self, doc_key: str) -> str:
        """Retrieve document content from S3 (supports text, PDF, DOCX, CSV, Excel, PPT)"""
        try:
            # Skip binary files that can't be indexed
            skip_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.ico', '.svg', 
                              '.mp4', '.avi', '.mov', '.mp3', '.wav', '.zip', '.tar', '.gz']
            if any(doc_key.lower().endswith(ext) for ext in skip_extensions):
                logger.debug(f"Skipping non-indexable file: {doc_key}")
                return ""
            
            response = s3_client.get_object(
                Bucket=S3_DOCS_BUCKET,
                Key=doc_key
            )
            
            doc_lower = doc_key.lower()
            file_bytes = response['Body'].read()
            
            if doc_lower.endswith('.pdf'):
                return self._extract_pdf_text(file_bytes, doc_key)
            elif doc_lower.endswith('.docx'):
                return self._extract_docx_text(file_bytes, doc_key)
            elif doc_lower.endswith('.csv'):
                return self._extract_csv_text(file_bytes, doc_key)
            elif doc_lower.endswith(('.xlsx', '.xls')):
                return self._extract_excel_text(file_bytes, doc_key)
            elif doc_lower.endswith(('.pptx', '.ppt')):
                return self._extract_ppt_text(file_bytes, doc_key)
            else:
                # Try to decode as UTF-8, skip if it fails
                try:
                    content = file_bytes.decode('utf-8')
                    return content
                except UnicodeDecodeError:
                    logger.warning(f"Skipping binary file (UTF-8 decode failed): {doc_key}")
                    return ""
        except Exception as e:
            logger.error(f"Error reading document {doc_key}: {str(e)}")
            return ""
    
    def _extract_pdf_text(self, pdf_bytes: bytes, doc_key: str) -> str:
        """Extract text from PDF using PyPDF2"""
        try:
            import io
            try:
                from PyPDF2 import PdfReader
            except ImportError:
                logger.error("PyPDF2 not installed. Install with: pip install PyPDF2")
                return ""
            
            pdf_file = io.BytesIO(pdf_bytes)
            reader = PdfReader(pdf_file)
            
            text_parts = []
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
            
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(extracted_text)} chars from PDF: {doc_key} ({len(reader.pages)} pages)")
            return extracted_text
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            return ""
    
    def _extract_docx_text(self, docx_bytes: bytes, doc_key: str) -> str:
        """Extract text from DOCX using python-docx"""
        try:
            import io
            try:
                from docx import Document
            except ImportError:
                logger.error("python-docx not installed. Install with: pip install python-docx")
                return ""
            
            docx_file = io.BytesIO(docx_bytes)
            doc = Document(docx_file)
            
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(extracted_text)} chars from DOCX: {doc_key}")
            return extracted_text
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            return ""
    
    def _extract_csv_text(self, csv_bytes: bytes, doc_key: str) -> str:
        """Extract text from CSV"""
        try:
            import io
            import csv
            
            csv_file = io.StringIO(csv_bytes.decode('utf-8'))
            reader = csv.reader(csv_file)
            
            text_parts = []
            for row_num, row in enumerate(reader):
                if row:
                    text_parts.append(f"[Row {row_num + 1}] " + " | ".join(row))
            
            extracted_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(extracted_text)} chars from CSV: {doc_key} ({len(text_parts)} rows)")
            return extracted_text
        except Exception as e:
            logger.error(f"Error extracting CSV text: {str(e)}")
            return ""
    
    def _extract_excel_text(self, excel_bytes: bytes, doc_key: str) -> str:
        """Extract text from Excel (xlsx/xls) using openpyxl"""
        try:
            import io
            try:
                from openpyxl import load_workbook
            except ImportError:
                logger.error("openpyxl not installed. Install with: pip install openpyxl")
                return ""
            
            excel_file = io.BytesIO(excel_bytes)
            workbook = load_workbook(excel_file, read_only=True, data_only=True)
            
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                
                for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        text_parts.append(f"[Row {row_num}] " + " | ".join(row_values))
            
            workbook.close()
            extracted_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(extracted_text)} chars from Excel: {doc_key} ({len(workbook.sheetnames)} sheets)")
            return extracted_text
        except Exception as e:
            logger.error(f"Error extracting Excel text: {str(e)}")
            return ""
    
    def _extract_ppt_text(self, ppt_bytes: bytes, doc_key: str) -> str:
        """Extract text from PowerPoint (pptx) using python-pptx"""
        try:
            import io
            try:
                from pptx import Presentation
            except ImportError:
                logger.error("python-pptx not installed. Install with: pip install python-pptx")
                return ""
            
            ppt_file = io.BytesIO(ppt_bytes)
            prs = Presentation(ppt_file)
            
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))
            
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(extracted_text)} chars from PPT: {doc_key} ({len(prs.slides)} slides)")
            return extracted_text
        except Exception as e:
            logger.error(f"Error extracting PPT text: {str(e)}")
            return ""
    
    def index_document(self, doc_key: str, chunk_size: int = 1000) -> int:
        """Index a document by chunking and storing vectors"""
        if not self._ensure_vector_index():
            return 0
        
        content = self.get_document_content(doc_key)
        if not content:
            return 0
        
        chunks = []
        overlap = 200
        for i in range(0, len(content), chunk_size - overlap):
            chunk = content[i:i + chunk_size]
            if chunk.strip():
                chunks.append(chunk)
        
        vectors = []
        doc_name = doc_key.replace(self.docs_prefix, '')
        
        for idx, chunk in enumerate(chunks):
            try:
                embedding = self._get_embedding(chunk)
                vector_key = f"{doc_key}_chunk_{idx}"
                vectors.append({
                    'key': vector_key,
                    'data': {'float32': embedding},
                    'metadata': {
                        'doc_key': doc_key,
                        'doc_name': doc_name,
                        'chunk_index': idx,
                        'content': chunk[:500]
                    }
                })
            except Exception as e:
                logger.error(f"Error creating embedding for chunk {idx}: {str(e)}")
                continue
        
        if vectors:
            try:
                batch_size = 100
                for i in range(0, len(vectors), batch_size):
                    batch = vectors[i:i + batch_size]
                    s3vectors_client.put_vectors(
                        vectorBucketName=S3_VECTOR_BUCKET,
                        indexName=self.index_name,
                        vectors=batch
                    )
                logger.info(f"Indexed {len(vectors)} chunks from {doc_name}")
            except Exception as e:
                logger.error(f"Error storing vectors: {str(e)}")
                return 0
        
        return len(vectors)
    
    def index_all_documents(self) -> Dict[str, int]:
        """Index all documents in the project folder"""
        documents = self.list_documents()
        results = {}
        
        # Skip non-indexable files
        skip_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.ico', '.svg', 
                          '.mp4', '.avi', '.mov', '.mp3', '.wav', '.zip', '.tar', '.gz']
        
        for doc in documents:
            doc_key = doc['key']
            doc_name = doc['name']
            
            # Skip binary files
            if any(doc_name.lower().endswith(ext) for ext in skip_extensions):
                logger.debug(f"Skipping non-indexable file: {doc_name}")
                results[doc_name] = 0
                continue
            
            chunks_indexed = self.index_document(doc_key)
            results[doc_name] = chunks_indexed
        
        return results
    
    def query(self, query_text: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """Query the vector index for relevant document chunks"""
        if not self._ensure_vector_index():
            return []
        
        try:
            query_embedding = self._get_embedding(query_text, purpose="GENERIC_RETRIEVAL")
            
            response = s3vectors_client.query_vectors(
                vectorBucketName=S3_VECTOR_BUCKET,
                indexName=self.index_name,
                topK=top_k,
                queryVector={'float32': query_embedding},
                returnMetadata=True,
                returnDistance=True
            )
            
            results = []
            for vector in response.get('vectors', []):
                metadata = vector.get('metadata', {})
                results.append({
                    'key': vector.get('key'),
                    'distance': vector.get('distance'),
                    'doc_name': metadata.get('doc_name', ''),
                    'content': metadata.get('content', ''),
                    'chunk_index': metadata.get('chunk_index', 0)
                })
            
            logger.info(f"Query returned {len(results)} results")
            return results
        except Exception as e:
            logger.error(f"Error querying vectors: {str(e)}")
            return []
    
    def get_context_for_query(self, query_text: str, top_k: int = 5) -> str:
        """Get formatted context string for RAG"""
        results = self.query(query_text, top_k)
        
        if not results:
            return "No relevant documents found."
        
        context_parts = []
        for i, result in enumerate(results, 1):
            context_parts.append(
                f"[Source {i}: {result['doc_name']}]\n{result['content']}"
            )
        
        return "\n\n---\n\n".join(context_parts)


# =============================================================================
# WEB CRAWLER FOR KNOWLEDGE BASE
# =============================================================================
class WebCrawler:
    """
    Web crawler that fetches and indexes web pages into S3 Vectors.
    Uses BeautifulSoup4 for HTML parsing and respects crawling scope.
    """
    
    def __init__(self, rag_instance: S3VectorsRAG):
        self.rag = rag_instance
        self.visited_urls = set()
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (compatible; QubitzBot/1.0; +https://qubitz.ai)'
        })
    
    def _is_valid_url(self, url: str, seed_url: str, scope: str = 'HOST_ONLY') -> bool:
        """Check if URL should be crawled based on scope."""
        from urllib.parse import urlparse
        
        seed_parsed = urlparse(seed_url)
        url_parsed = urlparse(url)
        
        if scope == 'HOST_ONLY':
            return url_parsed.netloc == seed_parsed.netloc
        elif scope == 'SUBDOMAINS':
            seed_domain = '.'.join(seed_parsed.netloc.split('.')[-2:])
            url_domain = '.'.join(url_parsed.netloc.split('.')[-2:])
            return url_domain == seed_domain
        else:  # DEFAULT - same host only
            return url_parsed.netloc == seed_parsed.netloc
    
    def _extract_text_from_html(self, html: str) -> str:
        """Extract clean text from HTML using BeautifulSoup."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, 'lxml')
            
            # Remove script and style elements
            for script in soup(['script', 'style', 'nav', 'footer', 'header']):
                script.decompose()
            
            # Get text
            text = soup.get_text(separator='\n', strip=True)
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from HTML: {e}")
            return ""
    
    def crawl_url(self, url: str, max_depth: int = 2, current_depth: int = 0) -> Dict[str, Any]:
        """
        Crawl a single URL and extract content.
        
        Returns dict with:
            - url: The crawled URL
            - title: Page title
            - content: Extracted text content
            - links: List of discovered links
            - success: Boolean indicating if crawl succeeded
        """
        if url in self.visited_urls or current_depth > max_depth:
            return {'success': False, 'reason': 'already_visited_or_max_depth'}
        
        self.visited_urls.add(url)
        
        try:
            response = self.session.get(url, timeout=10, allow_redirects=True)
            response.raise_for_status()
            
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(response.content, 'lxml')
            
            # Extract title
            title = soup.title.string if soup.title else url
            
            # Extract text content
            content = self._extract_text_from_html(response.text)
            
            # Extract links for further crawling
            links = []
            for link in soup.find_all('a', href=True):
                href = link['href']
                from urllib.parse import urljoin
                absolute_url = urljoin(url, href)
                if absolute_url.startswith('http'):
                    links.append(absolute_url)
            
            logger.info(f"Crawled: {url} ({len(content)} chars, {len(links)} links)")
            
            return {
                'success': True,
                'url': url,
                'title': title,
                'content': content,
                'links': links,
                'size': len(content)
            }
            
        except Exception as e:
            logger.error(f"Error crawling {url}: {e}")
            return {'success': False, 'url': url, 'error': str(e)}
    
    def crawl_and_index(
        self,
        seed_urls: List[str],
        crawling_scope: str = 'HOST_ONLY',
        max_pages: int = 50,
        max_depth: int = 2
    ) -> Dict[str, Any]:
        """
        Crawl seed URLs and index content into S3 Vectors.
        
        Args:
            seed_urls: List of starting URLs
            crawling_scope: 'HOST_ONLY', 'SUBDOMAINS', or 'DEFAULT'
            max_pages: Maximum number of pages to crawl
            max_depth: Maximum depth to crawl from seed URLs
        
        Returns:
            Dict with crawl statistics and indexed pages
        """
        from collections import deque
        
        queue = deque([(url, 0) for url in seed_urls])
        crawled_pages = []
        indexed_count = 0
        failed_count = 0
        
        while queue and len(crawled_pages) < max_pages:
            url, depth = queue.popleft()
            
            # Check if URL is valid for crawling
            seed_url = seed_urls[0]  # Use first seed as reference
            if not self._is_valid_url(url, seed_url, crawling_scope):
                continue
            
            # Crawl the URL
            result = self.crawl_url(url, max_depth, depth)
            
            if result.get('success'):
                crawled_pages.append(result)
                
                # Index the content into S3 Vectors
                try:
                    # Create a temporary "document" in S3 format
                    doc_key = f"web_crawled/{url.replace('://', '_').replace('/', '_')[:200]}.txt"
                    
                    # Store content chunks directly in vector index
                    content = result['content']
                    if content:
                        chunk_size = 1000
                        overlap = 200
                        chunks = []
                        
                        for i in range(0, len(content), chunk_size - overlap):
                            chunk = content[i:i + chunk_size]
                            if chunk.strip():
                                chunks.append(chunk)
                        
                        # Create embeddings and store vectors
                        vectors = []
                        for idx, chunk in enumerate(chunks):
                            try:
                                embedding = self.rag._get_embedding(chunk)
                                vector_key = f"{doc_key}_chunk_{idx}"
                                vectors.append({
                                    'key': vector_key,
                                    'data': {'float32': embedding},
                                    'metadata': {
                                        'doc_key': doc_key,
                                        'doc_name': result['title'],
                                        'url': url,
                                        'chunk_index': idx,
                                        'content': chunk[:500]
                                    }
                                })
                            except Exception as e:
                                logger.error(f"Error creating embedding for {url} chunk {idx}: {e}")
                        
                        if vectors:
                            self.rag._ensure_vector_index()
                            batch_size = 100
                            for i in range(0, len(vectors), batch_size):
                                batch = vectors[i:i + batch_size]
                                s3vectors_client.put_vectors(
                                    vectorBucketName=S3_VECTOR_BUCKET,
                                    indexName=self.rag.index_name,
                                    vectors=batch
                                )
                            indexed_count += 1
                            logger.info(f"Indexed {len(vectors)} chunks from {url}")
                
                except Exception as e:
                    logger.error(f"Error indexing {url}: {e}")
                    failed_count += 1
                
                # Add discovered links to queue
                if depth < max_depth:
                    for link in result.get('links', []):
                        if link not in self.visited_urls:
                            queue.append((link, depth + 1))
            else:
                failed_count += 1
        
        return {
            'crawled_pages': len(crawled_pages),
            'indexed_pages': indexed_count,
            'failed_pages': failed_count,
            'visited_urls': list(self.visited_urls),
            'pages': crawled_pages
        }


# Global RAG instance
rag_instance = None


def get_rag_instance() -> S3VectorsRAG:
    """Get or create RAG instance"""
    global rag_instance
    if rag_instance is None:
        rag_instance = S3VectorsRAG(USER_ID, PROJECT_ID)
    return rag_instance


# =============================================================================
# BEDROCK KNOWLEDGE BASE RAG (per user_id/project_id)
# =============================================================================
# NOTE: This class exists for backward compatibility but is NOT used by default.
# The system uses CUSTOM DATA SOURCE (S3 Vectors) for all connectors instead.
# See get_kb_rag_instance() which returns S3VectorsRAG, not BedrockKnowledgeBaseRAG.
KNOWLEDGE_BASE_ID = os.environ.get('KNOWLEDGE_BASE_ID', '')
KB_REGION = os.environ.get('KB_REGION', AWS_REGION)


class BedrockKnowledgeBaseRAG:
    """
    RAG implementation using Bedrock Knowledge Base.
    
    IMPORTANT: This class is NOT used by default. The system uses S3VectorsRAG (custom data source)
    for all connectors (S3, web crawler, Confluence, SharePoint, Salesforce).
    
    This class is only used if KNOWLEDGE_BASE_ID environment variable is explicitly set.
    Queries a provisioned KB that is scoped per user_id/project_id.
    Supports: S3, web_crawler, Confluence, SharePoint, Salesforce data sources.
    """

    def __init__(self, kb_id: str, region: str = 'eu-central-1'):
        self.kb_id = kb_id
        self.region = region
        self._runtime_client = boto3.client('bedrock-agent-runtime', region_name=region)
        self._agent_client = boto3.client('bedrock-agent', region_name=region)
        logger.info(f"BedrockKnowledgeBaseRAG initialized: kb_id={kb_id}, region={region}")

    def query(self, query_text: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """Query the knowledge base for relevant content."""
        if not self.kb_id:
            return []

        try:
            response = self._runtime_client.retrieve(
                knowledgeBaseId=self.kb_id,
                retrievalQuery={'text': query_text},
                retrievalConfiguration={
                    'vectorSearchConfiguration': {
                        'numberOfResults': top_k
                    }
                }
            )

            results = []
            for item in response.get('retrievalResults', []):
                content = item.get('content', {}).get('text', '')
                location = item.get('location', {})
                score = item.get('score', 0.0)
                metadata = item.get('metadata', {})

                # Extract source info based on location type
                source_type = location.get('type', 'UNKNOWN')
                source_uri = ''
                if source_type == 'S3':
                    source_uri = location.get('s3Location', {}).get('uri', '')
                elif source_type == 'WEB':
                    source_uri = location.get('webLocation', {}).get('url', '')
                elif source_type == 'CONFLUENCE':
                    source_uri = location.get('confluenceLocation', {}).get('url', '')
                elif source_type == 'SHAREPOINT':
                    source_uri = location.get('sharePointLocation', {}).get('url', '')
                elif source_type == 'SALESFORCE':
                    source_uri = location.get('salesforceLocation', {}).get('url', '')

                results.append({
                    'content': content,
                    'source_type': source_type,
                    'source_uri': source_uri,
                    'score': score,
                    'metadata': metadata
                })

            logger.info(f"KB query returned {len(results)} results for: {query_text[:80]}...")
            return results

        except Exception as e:
            logger.error(f"KB query error: {str(e)}")
            return []

    def get_context_for_query(self, query_text: str, top_k: int = 5) -> str:
        """Get formatted context string from KB for RAG augmentation."""
        results = self.query(query_text, top_k)

        if not results:
            return "No relevant information found in the knowledge base."

        context_parts = []
        for i, result in enumerate(results, 1):
            source_info = f"{result['source_type']}"
            if result['source_uri']:
                source_info += f": {result['source_uri']}"
            context_parts.append(
                f"[Source {i} ({source_info}, score: {result['score']:.3f})]\n{result['content']}"
            )

        return "\n\n---\n\n".join(context_parts)

    def list_data_sources(self) -> List[Dict[str, Any]]:
        """List all data sources connected to this knowledge base."""
        if not self.kb_id:
            return []

        try:
            response = self._agent_client.list_data_sources(
                knowledgeBaseId=self.kb_id
            )

            sources = []
            for ds in response.get('dataSourceSummaries', []):
                sources.append({
                    'data_source_id': ds.get('dataSourceId', ''),
                    'name': ds.get('name', ''),
                    'status': ds.get('status', ''),
                    'updated_at': ds.get('updatedAt', '').isoformat() if ds.get('updatedAt') else ''
                })

            return sources

        except Exception as e:
            logger.error(f"Error listing KB data sources: {str(e)}")
            return []


# Global KB RAG instance - uses CUSTOM DATA SOURCE (S3 Vectors) for all connectors
_kb_rag_instance: Optional[S3VectorsRAG] = None
_bedrock_kb_rag_instance: Optional[BedrockKnowledgeBaseRAG] = None


def get_kb_rag_instance() -> Optional[S3VectorsRAG]:
    """
    Get or create Knowledge Base RAG instance using CUSTOM DATA SOURCE.
    
    IMPORTANT: This uses S3VectorsRAG (custom S3 Vectors implementation) for ALL data sources:
    - S3 buckets
    - Web crawler results
    - Confluence pages (crawled and indexed)
    - SharePoint documents (crawled and indexed)
    - Salesforce records (crawled and indexed)
    
    This provides full control over indexing and retrieval without Bedrock Knowledge Base.
    """
    global _kb_rag_instance
    if _kb_rag_instance is None:
        # Initialize with user_id/project_id scoping
        _kb_rag_instance = S3VectorsRAG(USER_ID, PROJECT_ID)
        logger.info("Initialized custom data source RAG (S3 Vectors) for all connectors")
    return _kb_rag_instance


def get_bedrock_kb_rag_instance() -> Optional[BedrockKnowledgeBaseRAG]:
    """
    Get or create Bedrock Knowledge Base RAG instance.
    
    NOTE: This is ONLY used if KNOWLEDGE_BASE_ID is explicitly set.
    By default, we use custom data sources (S3 Vectors) instead.
    """
    global _bedrock_kb_rag_instance
    if _bedrock_kb_rag_instance is None and KNOWLEDGE_BASE_ID:
        _bedrock_kb_rag_instance = BedrockKnowledgeBaseRAG(KNOWLEDGE_BASE_ID, KB_REGION)
        logger.info(f"Initialized Bedrock Knowledge Base RAG: {KNOWLEDGE_BASE_ID}")
    return _bedrock_kb_rag_instance


@tool
def knowledge_base_query(query: str, top_k: int = 5) -> str:
    """
    Query the Knowledge Base for relevant information from connected data sources.

    USES CUSTOM DATA SOURCE (S3 Vectors) for all connectors:
    - S3 buckets: Documents indexed from configured S3 paths
    - Web crawler: Pages crawled and indexed into S3 Vectors
    - Confluence: Pages fetched via API and indexed into S3 Vectors
    - SharePoint: Documents fetched via API and indexed into S3 Vectors
    - Salesforce: Records fetched via API and indexed into S3 Vectors

    This provides unified search across ALL data sources without Bedrock Knowledge Base.

    Args:
        query: The search query to find relevant content
        top_k: Number of top results to return (default: 5)

    Returns:
        Relevant content excerpts with source attribution
    """
    # Always use custom data source (S3 Vectors) for all connectors
    kb_rag = get_kb_rag_instance()
    if not kb_rag:
        return "Knowledge Base not configured."

    try:
        context = kb_rag.get_context_for_query(query, top_k)
        logger.info(f"Custom data source query completed: {len(context)} chars returned")
        return context
    except Exception as e:
        logger.error(f"Error in knowledge_base_query (custom data source): {str(e)}")
        return f"Error querying knowledge base: {str(e)}"


@tool
def knowledge_base_list_sources() -> str:
    """
    List all documents indexed in the knowledge base (custom data source).

    Returns documents from ALL connector types:
    - S3 buckets
    - Web crawler results
    - Confluence pages
    - SharePoint documents
    - Salesforce records

    Returns:
        List of indexed documents with metadata
    """
    kb_rag = get_kb_rag_instance()
    if not kb_rag:
        return "Knowledge Base not configured."

    try:
        documents = kb_rag.list_documents()
        if not documents:
            return f"No documents found in knowledge base for {USER_ID}/{PROJECT_ID}."

        lines = [f"Knowledge Base ({USER_ID}/{PROJECT_ID}) - {len(documents)} document(s):"]
        for doc in documents:
            lines.append(
                f"  - {doc['name']} ({doc['size']} bytes, "
                f"modified: {doc['last_modified']})"
            )
        return "\n".join(lines)

    except Exception as e:
        logger.error(f"Error listing KB sources: {str(e)}")
        return f"Error listing knowledge base sources: {str(e)}"


@tool
def crawl_and_index_website(
    url: str,
    max_pages: int = 20,
    crawling_scope: str = 'HOST_ONLY'
) -> str:
    """
    Crawl a website and index its content into the knowledge base.
    
    This tool fetches web pages, extracts text content, and indexes it
    into S3 Vectors for later retrieval via knowledge_base_query.
    
    Args:
        url: The starting URL to crawl (seed URL)
        max_pages: Maximum number of pages to crawl (default: 20, max: 100)
        crawling_scope: Crawling scope - 'HOST_ONLY' (same host), 'SUBDOMAINS' (include subdomains), or 'DEFAULT'
    
    Returns:
        Summary of crawled and indexed pages
    
    Example:
        crawl_and_index_website("https://docs.example.com", max_pages=50, crawling_scope="HOST_ONLY")
    """
    try:
        # Limit max_pages to prevent abuse
        max_pages = min(max_pages, 100)
        
        rag = get_rag_instance()
        crawler = WebCrawler(rag)
        
        logger.info(f"Starting web crawl: {url} (max_pages={max_pages}, scope={crawling_scope})")
        
        result = crawler.crawl_and_index(
            seed_urls=[url],
            crawling_scope=crawling_scope,
            max_pages=max_pages,
            max_depth=2
        )
        
        summary = [
            f"Web Crawl Complete:",
            f"  - Crawled: {result['crawled_pages']} pages",
            f"  - Indexed: {result['indexed_pages']} pages",
            f"  - Failed: {result['failed_pages']} pages",
            f"  - Seed URL: {url}",
            f"  - Scope: {crawling_scope}",
            "",
            "Content is now searchable via knowledge_base_query tool."
        ]
        
        return "\n".join(summary)
        
    except Exception as e:
        logger.error(f"Error in crawl_and_index_website: {e}")
        return f"Error crawling website: {str(e)}"


# =============================================================================
# CONNECTOR INDEXING TOOLS (Confluence, SharePoint, Salesforce)
# =============================================================================
# These tools fetch content from external systems and index into S3 Vectors
# using CUSTOM DATA SOURCE approach (not Bedrock Knowledge Base)

def index_confluence_pages(
    space_keys: List[str],
    max_pages: int = 100
) -> str:
    """
    Fetch and index Confluence pages into the knowledge base using MCP server via AgentCore Gateway.

    This tool connects to Confluence through the MCP server, fetches pages from specified spaces,
    and indexes them into S3 Vectors for semantic search.

    Prerequisites:
    - Confluence MCP server must be configured via AgentCore Gateway
    - Credentials are resolved dynamically per user_id + project_id from DynamoDB ToolCredentials

    Args:
        space_keys: List of Confluence space keys to index (e.g., ['DOCS', 'TEAM'])
        max_pages: Maximum number of pages to index per space (default: 100)

    Returns:
        Summary of indexed Confluence pages
    """
    try:
        rag = get_kb_rag_instance()
        if not rag:
            return "Knowledge Base not configured."

        gateway = get_mcp_gateway()
        if not gateway or not gateway.is_connected:
            return (
                "AgentCore Gateway not available. Please configure Confluence MCP server via Gateway.\n"
                "Use set_tool action with tool='rag_mcp_server' and connector_type='confluence' to configure the Confluence connector."
            )

        logger.info(f"Indexing Confluence pages via MCP, spaces: {space_keys}")

        # Find Confluence MCP tools from gateway
        confluence_tools = [t for t in gateway.gateway_tools
                           if 'confluence' in getattr(t, 'tool_name', '').lower()]

        if not confluence_tools:
            return (
                "Confluence MCP server not found in gateway. Available tools: " +
                ", ".join([getattr(t, 'tool_name', str(t)) for t in gateway.gateway_tools])
            )

        indexed_count = 0
        failed_count = 0
        indexed_pages = []

        # Fetch pages from each space via MCP
        for space_key in space_keys:
            try:
                # Try to find list_pages or search_pages tool
                list_tool = None
                for tool in confluence_tools:
                    tool_name = getattr(tool, 'tool_name', '')
                    if 'list' in tool_name.lower() or 'search' in tool_name.lower() or 'get' in tool_name.lower():
                        list_tool = tool
                        break

                if not list_tool:
                    logger.warning(f"No list/search tool found for Confluence, using first available tool")
                    list_tool = confluence_tools[0]

                # Call MCP tool to fetch pages — include user_id/project_id for credential resolution
                tool_name = getattr(list_tool, 'tool_name', str(list_tool))
                logger.info(f"Calling Confluence MCP tool: {tool_name} for space {space_key}")

                # Base params always include user_id + project_id for dynamic credential lookup
                base_params = {"user_id": USER_ID, "project_id": PROJECT_ID}

                # Try different parameter formats based on common MCP patterns
                try:
                    result = gateway.mcp_client.call_tool_sync(
                        tool_name,
                        {**base_params, "space_key": space_key, "limit": max_pages}
                    )
                except Exception as e1:
                    try:
                        result = gateway.mcp_client.call_tool_sync(
                            tool_name,
                            {**base_params, "spaceKey": space_key, "limit": max_pages}
                        )
                    except Exception as e2:
                        result = gateway.mcp_client.call_tool_sync(
                            tool_name,
                            {**base_params, "space": space_key, "max_results": max_pages}
                        )

                # Parse result and extract page content
                if result:
                    result_str = str(result)
                    # Index the content into S3 Vectors
                    chunks = []
                    chunk_size = 1000
                    overlap = 200
                    for i in range(0, len(result_str), chunk_size - overlap):
                        chunk = result_str[i:i + chunk_size]
                        if chunk.strip():
                            chunks.append(chunk)

                    # Create embeddings and store vectors
                    for idx, chunk in enumerate(chunks):
                        try:
                            embedding = rag._get_embedding(chunk)
                            vector_key = f"confluence/{space_key}/page_{indexed_count}_chunk_{idx}"
                            rag._ensure_vector_index()
                            s3vectors_client.put_vectors(
                                vectorBucketName=S3_VECTOR_BUCKET,
                                indexName=rag.index_name,
                                vectors=[{
                                    'key': vector_key,
                                    'data': {'float32': embedding},
                                    'metadata': {
                                        'source': 'confluence',
                                        'space_key': space_key,
                                        'chunk_index': idx,
                                        'content': chunk[:500]
                                    }
                                }]
                            )
                        except Exception as e:
                            logger.error(f"Error creating embedding for Confluence chunk: {e}")
                            failed_count += 1

                    indexed_count += len(chunks)
                    indexed_pages.append(f"{space_key}: {len(chunks)} chunks")
                    logger.info(f"Indexed {len(chunks)} chunks from Confluence space {space_key}")

            except Exception as e:
                logger.error(f"Error fetching Confluence space {space_key}: {e}")
                failed_count += 1

        summary = [
            f"Confluence Indexing Complete (via MCP):",
            f"  - Spaces processed: {len(space_keys)}",
            f"  - Total chunks indexed: {indexed_count}",
            f"  - Failed: {failed_count}",
            "",
            "Indexed spaces:"
        ]
        summary.extend([f"  - {page}" for page in indexed_pages])
        summary.append("")
        summary.append("Content is now searchable via knowledge_base_query tool.")

        return "\n".join(summary)

    except Exception as e:
        logger.error(f"Error indexing Confluence pages: {str(e)}")
        return f"Error indexing Confluence: {str(e)}"


@tool
def index_sharepoint_documents(
    site_urls: List[str],
    max_documents: int = 100
) -> str:
    """
    Fetch and index SharePoint documents into the knowledge base using MCP server via AgentCore Gateway.
    
    This tool connects to SharePoint Online through the MCP server, fetches documents from specified sites,
    and indexes them into S3 Vectors for semantic search.
    
    Prerequisites:
    - SharePoint MCP server must be configured via AgentCore Gateway
    - Connection credentials (domain, tenant_id, client_id, client_secret) are configured at the gateway level
    
    Args:
        site_urls: List of SharePoint site URLs to index
        max_documents: Maximum number of documents to index per site (default: 100)
    
    Returns:
        Summary of indexed SharePoint documents
    """
    try:
        rag = get_kb_rag_instance()
        if not rag:
            return "Knowledge Base not configured."
        
        gateway = get_mcp_gateway()
        if not gateway or not gateway.is_connected:
            return (
                "AgentCore Gateway not available. Please configure SharePoint MCP server via Gateway.\n"
                "Use set_tool action with tool='rag_mcp_server' and connector_type='sharepoint' to configure the SharePoint connector."
            )
        
        logger.info(f"Indexing SharePoint documents via MCP, sites: {site_urls}")
        
        # Find SharePoint MCP tools from gateway
        sharepoint_tools = [t for t in gateway.gateway_tools 
                           if 'sharepoint' in getattr(t, 'tool_name', '').lower()]
        
        if not sharepoint_tools:
            return (
                "SharePoint MCP server not found in gateway. Available tools: " +
                ", ".join([getattr(t, 'tool_name', str(t)) for t in gateway.gateway_tools])
            )
        
        indexed_count = 0
        failed_count = 0
        indexed_docs = []
        
        # Fetch documents from each site via MCP
        for site_url in site_urls:
            try:
                # Try to find list_documents or get_documents tool
                list_tool = None
                for tool in sharepoint_tools:
                    tool_name = getattr(tool, 'tool_name', '')
                    if 'list' in tool_name.lower() or 'get' in tool_name.lower() or 'search' in tool_name.lower():
                        list_tool = tool
                        break
                
                if not list_tool:
                    logger.warning(f"No list/get tool found for SharePoint, using first available tool")
                    list_tool = sharepoint_tools[0]
                
                # Call MCP tool to fetch documents
                tool_name = getattr(list_tool, 'tool_name', str(list_tool))
                logger.info(f"Calling SharePoint MCP tool: {tool_name} for site {site_url}")
                
                # Base params include user_id + project_id for dynamic credential lookup
                base_params = {"user_id": USER_ID, "project_id": PROJECT_ID}
                
                # Try different parameter formats based on common MCP patterns
                try:
                    result = gateway.mcp_client.call_tool_sync(
                        tool_name,
                        {**base_params, "site_url": site_url, "limit": max_documents}
                    )
                except Exception as e1:
                    try:
                        result = gateway.mcp_client.call_tool_sync(
                            tool_name,
                            {**base_params, "siteUrl": site_url, "max_results": max_documents}
                        )
                    except Exception as e2:
                        result = gateway.mcp_client.call_tool_sync(
                            tool_name,
                            {**base_params, "site": site_url, "limit": max_documents}
                        )
                
                # Parse result and extract document content
                if result:
                    result_str = str(result)
                    # Index the content into S3 Vectors
                    chunks = []
                    chunk_size = 1000
                    overlap = 200
                    for i in range(0, len(result_str), chunk_size - overlap):
                        chunk = result_str[i:i + chunk_size]
                        if chunk.strip():
                            chunks.append(chunk)
                    
                    # Create embeddings and store vectors
                    for idx, chunk in enumerate(chunks):
                        try:
                            embedding = rag._get_embedding(chunk)
                            vector_key = f"sharepoint/{site_url.replace('/', '_')}/doc_{indexed_count}_chunk_{idx}"
                            rag._ensure_vector_index()
                            s3vectors_client.put_vectors(
                                vectorBucketName=S3_VECTOR_BUCKET,
                                indexName=rag.index_name,
                                vectors=[{
                                    'key': vector_key,
                                    'data': {'float32': embedding},
                                    'metadata': {
                                        'source': 'sharepoint',
                                        'site_url': site_url,
                                        'chunk_index': idx,
                                        'content': chunk[:500]
                                    }
                                }]
                            )
                        except Exception as e:
                            logger.error(f"Error creating embedding for SharePoint chunk: {e}")
                            failed_count += 1
                    
                    indexed_count += len(chunks)
                    indexed_docs.append(f"{site_url}: {len(chunks)} chunks")
                    logger.info(f"Indexed {len(chunks)} chunks from SharePoint site {site_url}")
                
            except Exception as e:
                logger.error(f"Error fetching SharePoint site {site_url}: {e}")
                failed_count += 1
        
        summary = [
            f"SharePoint Indexing Complete (via MCP):",
            f"  - Sites processed: {len(site_urls)}",
            f"  - Total chunks indexed: {indexed_count}",
            f"  - Failed: {failed_count}",
            "",
            "Indexed sites:"
        ]
        summary.extend([f"  - {doc}" for doc in indexed_docs])
        summary.append("")
        summary.append("Content is now searchable via knowledge_base_query tool.")
        
        return "\n".join(summary)
        
    except Exception as e:
        logger.error(f"Error indexing SharePoint documents: {str(e)}")
        return f"Error indexing SharePoint: {str(e)}"


@tool
def index_salesforce_records(
    object_types: List[str],
    max_records: int = 1000
) -> str:
    """
    Fetch and index Salesforce records into the knowledge base using MCP server via AgentCore Gateway.
    
    This tool connects to Salesforce through the MCP server, fetches records from specified object types,
    and indexes them into S3 Vectors for semantic search.
    
    Prerequisites:
    - Salesforce MCP server must be configured via AgentCore Gateway
    - Connection credentials (host_url, consumer_key, consumer_secret) are configured at the gateway level
    
    Args:
        object_types: List of Salesforce object types to index (e.g., ['Account', 'Case', 'Knowledge__kav'])
        max_records: Maximum number of records to index per object type (default: 1000)
    
    Returns:
        Summary of indexed Salesforce records
    """
    try:
        rag = get_kb_rag_instance()
        if not rag:
            return "Knowledge Base not configured."
        
        gateway = get_mcp_gateway()
        if not gateway or not gateway.is_connected:
            return (
                "AgentCore Gateway not available. Please configure Salesforce MCP server via Gateway.\n"
                "Use set_tool action with tool='rag_mcp_server' and connector_type='salesforce' to configure the Salesforce connector."
            )
        
        logger.info(f"Indexing Salesforce records via MCP, objects: {object_types}")
        
        # Find Salesforce MCP tools from gateway
        salesforce_tools = [t for t in gateway.gateway_tools 
                           if 'salesforce' in getattr(t, 'tool_name', '').lower() or 
                              'sfdc' in getattr(t, 'tool_name', '').lower()]
        
        if not salesforce_tools:
            return (
                "Salesforce MCP server not found in gateway. Available tools: " +
                ", ".join([getattr(t, 'tool_name', str(t)) for t in gateway.gateway_tools])
            )
        
        indexed_count = 0
        failed_count = 0
        indexed_objects = []
        
        # Fetch records from each object type via MCP
        for object_type in object_types:
            try:
                # Try to find query or search tool
                query_tool = None
                for tool in salesforce_tools:
                    tool_name = getattr(tool, 'tool_name', '')
                    if 'query' in tool_name.lower() or 'search' in tool_name.lower() or 'get' in tool_name.lower():
                        query_tool = tool
                        break
                
                if not query_tool:
                    logger.warning(f"No query/search tool found for Salesforce, using first available tool")
                    query_tool = salesforce_tools[0]
                
                # Call MCP tool to fetch records
                tool_name = getattr(query_tool, 'tool_name', str(query_tool))
                logger.info(f"Calling Salesforce MCP tool: {tool_name} for object {object_type}")
                
                # Base params include user_id + project_id for dynamic credential lookup
                base_params = {"user_id": USER_ID, "project_id": PROJECT_ID}
                
                # Try different parameter formats based on common MCP patterns
                try:
                    result = gateway.mcp_client.call_tool_sync(
                        tool_name,
                        {**base_params, "object_type": object_type, "limit": max_records}
                    )
                except Exception as e1:
                    try:
                        result = gateway.mcp_client.call_tool_sync(
                            tool_name,
                            {**base_params, "sobject": object_type, "max_results": max_records}
                        )
                    except Exception as e2:
                        # Try SOQL query format
                        soql = f"SELECT Id, Name FROM {object_type} LIMIT {max_records}"
                        result = gateway.mcp_client.call_tool_sync(
                            tool_name,
                            {**base_params, "query": soql}
                        )
                
                # Parse result and extract record content
                if result:
                    result_str = str(result)
                    # Index the content into S3 Vectors
                    chunks = []
                    chunk_size = 1000
                    overlap = 200
                    for i in range(0, len(result_str), chunk_size - overlap):
                        chunk = result_str[i:i + chunk_size]
                        if chunk.strip():
                            chunks.append(chunk)
                    
                    # Create embeddings and store vectors
                    for idx, chunk in enumerate(chunks):
                        try:
                            embedding = rag._get_embedding(chunk)
                            vector_key = f"salesforce/{object_type}/record_{indexed_count}_chunk_{idx}"
                            rag._ensure_vector_index()
                            s3vectors_client.put_vectors(
                                vectorBucketName=S3_VECTOR_BUCKET,
                                indexName=rag.index_name,
                                vectors=[{
                                    'key': vector_key,
                                    'data': {'float32': embedding},
                                    'metadata': {
                                        'source': 'salesforce',
                                        'object_type': object_type,
                                        'chunk_index': idx,
                                        'content': chunk[:500]
                                    }
                                }]
                            )
                        except Exception as e:
                            logger.error(f"Error creating embedding for Salesforce chunk: {e}")
                            failed_count += 1
                    
                    indexed_count += len(chunks)
                    indexed_objects.append(f"{object_type}: {len(chunks)} chunks")
                    logger.info(f"Indexed {len(chunks)} chunks from Salesforce object {object_type}")
                
            except Exception as e:
                logger.error(f"Error fetching Salesforce object {object_type}: {e}")
                failed_count += 1
        
        summary = [
            f"Salesforce Indexing Complete (via MCP):",
            f"  - Object types processed: {len(object_types)}",
            f"  - Total chunks indexed: {indexed_count}",
            f"  - Failed: {failed_count}",
            "",
            "Indexed objects:"
        ]
        summary.extend([f"  - {obj}" for obj in indexed_objects])
        summary.append("")
        summary.append("Content is now searchable via knowledge_base_query tool.")
        
        return "\n".join(summary)
        
    except Exception as e:
        logger.error(f"Error indexing Salesforce records: {str(e)}")
        return f"Error indexing Salesforce: {str(e)}"
        logger.error(f"Error indexing Salesforce records: {str(e)}")
        return f"Error indexing Salesforce: {str(e)}"


# =============================================================================
# GATEWAY TARGETS CONFIG (which targets exist and their types)
# =============================================================================
# Set by CDK deployment as JSON env var. Maps target names to types.
# Example: [{"name": "TavilySearch", "type": "web_search"}, {"name": "slack", "type": "mcp_server"}]
GATEWAY_TARGETS_JSON = os.environ.get('GATEWAY_TARGETS', '[]')
try:
    GATEWAY_TARGETS = json.loads(GATEWAY_TARGETS_JSON) if GATEWAY_TARGETS_JSON else []
except (json.JSONDecodeError, TypeError):
    GATEWAY_TARGETS = []

# Build lookup: target_name → type
_TARGET_TYPE_MAP = {t.get('name', ''): t.get('type', '') for t in GATEWAY_TARGETS if t.get('name')}

# Known web search tool names from MCP servers (used for filtering)
# Tavily MCP exposes: tavily_search, tavily_extract, tavily_map, tavily_crawl, tavily_research
# SerpAPI MCP exposes: search
_WEB_SEARCH_TOOL_NAMES = {'tavily_search', 'tavily_extract', 'tavily_map', 'tavily_crawl', 'tavily_research', 'search'}


# =============================================================================
# MCP GATEWAY CLIENT (connects to AgentCore Gateway for all targets)
# =============================================================================
class MCPGatewayManager:
    """
    Manages MCP client connection to AgentCore Gateway.
    All targets (Tavily, SerpAPI, custom MCP servers) are auto-discovered.
    Per-agent tool filtering ensures each agent only sees its assigned targets.
    """

    def __init__(self, gateway_url: str):
        self.gateway_url = gateway_url
        self.mcp_client = None
        self.gateway_tools = []
        self._connected = False
        # Maps tool_name → target_name for filtering
        self._tool_to_target = {}

    def connect(self) -> List:
        """Connect to gateway and discover available tools.

        Uses IAM SigV4 authentication (via mcp_proxy_for_aws) when available,
        falls back to plain streamablehttp_client for non-IAM gateways.
        """
        if self._connected:
            return self.gateway_tools

        try:
            from strands.tools.mcp.mcp_client import MCPClient

            logger.info(f"Connecting to gateway MCP endpoint: {self.gateway_url}")

            # Try IAM SigV4 transport first (same pattern as github-mcp.py)
            try:
                from mcp_proxy_for_aws.client import aws_iam_streamablehttp_client
                logger.info("Using IAM SigV4 authentication for gateway")
                self.mcp_client = MCPClient(
                    lambda: aws_iam_streamablehttp_client(
                        endpoint=self.gateway_url,
                        aws_region=AWS_REGION,
                        aws_service="bedrock-agentcore"
                    )
                )
            except ImportError:
                from mcp.client.streamable_http import streamablehttp_client
                logger.info("mcp_proxy_for_aws not available, using plain HTTP transport")
                self.mcp_client = MCPClient(
                    lambda: streamablehttp_client(self.gateway_url)
                )

            self.mcp_client.__enter__()

            # Discover all tools from gateway with pagination
            all_tools = []
            more_tools = True
            pagination_token = None

            while more_tools:
                batch = self.mcp_client.list_tools_sync(pagination_token=pagination_token)
                all_tools.extend(batch)
                if hasattr(batch, 'pagination_token') and batch.pagination_token:
                    pagination_token = batch.pagination_token
                else:
                    more_tools = False

            self.gateway_tools = all_tools
            self._connected = True

            tool_names = [getattr(t, 'tool_name', str(t)) for t in self.gateway_tools]
            logger.info(f"Gateway connected: {len(self.gateway_tools)} tool(s) discovered: {tool_names}")

            return self.gateway_tools

        except ImportError:
            logger.warning("MCP client not available (strands.tools.mcp not installed)")
            return []
        except Exception as e:
            logger.error(f"Failed to connect to gateway: {str(e)}")
            return []

    @staticmethod
    def _is_web_search_tool(tool) -> bool:
        """Check if a tool is a web search tool (from Tavily or SerpAPI MCP)."""
        name = getattr(tool, 'tool_name', str(tool))
        # Handle gateway prefix: TargetName___ToolName
        base_name = name.split('___')[-1] if '___' in name else name
        return base_name in _WEB_SEARCH_TOOL_NAMES

    def get_tools_for_target(self, target_name: str) -> List:
        """Get tools belonging to a specific gateway target.

        Filtering strategy:
          1. Known web search tools are matched by name (tavily_search, tavily_extract, search)
          2. For other targets, all non-web-search tools are returned
             (gateway doesn't expose target→tool mapping via MCP, so we use heuristics)
        """
        if not self._connected:
            self.connect()

        target_type = _TARGET_TYPE_MAP.get(target_name, '')

        if target_type == 'web_search' or target_name in ('TavilySearch', 'SerpApiSearch'):
            # Return only web search tools
            return [t for t in self.gateway_tools if self._is_web_search_tool(t)]

        # For MCP server targets, return tools that are NOT web search
        # This works because web search tools have known names, and everything else
        # comes from MCP server targets
        return [t for t in self.gateway_tools if not self._is_web_search_tool(t)]

    def get_all_non_websearch_tools(self) -> List:
        """Get all tools that are NOT web search (Tavily/SerpAPI)."""
        if not self._connected:
            self.connect()
        return [t for t in self.gateway_tools if not self._is_web_search_tool(t)]

    def disconnect(self):
        """Disconnect from gateway."""
        if self.mcp_client and self._connected:
            try:
                self.mcp_client.__exit__(None, None, None)
                logger.info("Gateway MCP client disconnected")
            except Exception as e:
                logger.warning(f"Error disconnecting MCP client: {str(e)}")
            self._connected = False
            self.gateway_tools = []

    @property
    def is_connected(self) -> bool:
        return self._connected


# Global MCP gateway manager
_mcp_gateway: Optional[MCPGatewayManager] = None


def get_mcp_gateway() -> Optional[MCPGatewayManager]:
    """Get or create MCP gateway connection (singleton)."""
    global _mcp_gateway
    if _mcp_gateway is None and GATEWAY_URL:
        _mcp_gateway = MCPGatewayManager(GATEWAY_URL)
        _mcp_gateway.connect()
    return _mcp_gateway


# =============================================================================
# =============================================================================
# WEB SEARCH TOOL — routes to selected provider (playwright_scraper or gateway)
# =============================================================================

# Provider set via env var during deployment (from set_tool)
WEB_SEARCH_PROVIDER = os.environ.get('WEB_SEARCH_PROVIDER', 'playwright_scraper')


def _search_playwright(query: str, k: int, include_images: bool) -> str:
    """Playwright Scraper — direct HTTP, no API key needed."""
    import requests

    PRIMARY_API = "https://9bn49m5pk4.execute-api.us-east-1.amazonaws.com"
    QUICK_API = "https://77w5isieui.execute-api.us-east-1.amazonaws.com"

    payload = {"query": query, "k": k}
    if include_images:
        payload["include_images"] = True

    def make_request(api_url: str, api_name: str) -> dict:
        try:
            logger.info(f"Attempting web search via {api_name}: query='{query}', k={k}")
            response = requests.post(api_url, json=payload, headers={"Content-Type": "application/json"}, timeout=60)
            if response.status_code == 200:
                logger.info(f"{api_name} search successful")
                return response.json()
            else:
                return {"error": f"{api_name} returned status {response.status_code}"}
        except requests.Timeout:
            return {"error": f"{api_name} request timed out"}
        except Exception as e:
            return {"error": f"{api_name} error: {str(e)}"}

    result = make_request(QUICK_API, "Quick API")
    if "error" in result or not result.get("results"):
        logger.info("Quick API failed, falling back to Primary API")
        result = make_request(PRIMARY_API, "Primary API")
        if "error" in result:
            return f"Web search failed: {result.get('error', 'Unknown error')}"

    return _format_playwright_results(result, query, include_images)


def _search_via_gateway(query: str, k: int, provider: str) -> str:
    """Search via AgentCore Gateway (Tavily or SerpAPI MCP targets).

    The gateway handles authentication via credential providers.
    We call the gateway's MCP endpoint which routes to the correct target.
    """
    gateway = get_mcp_gateway()
    if not gateway or not gateway.is_connected:
        logger.error(f"Gateway not available for {provider}")
        return f"Web search failed: AgentCore Gateway not available for {provider}. Ensure gateway is configured and deployed."

    try:
        # Find the right tool from gateway's discovered tools
        # Gateway may prefix tool names: TargetName___ToolName (e.g. TavilySearch___tavily_search)
        tool_name_map = {
            'tavily': 'tavily_search',
            'serpapi': 'search',
        }
        target_tool_name = tool_name_map.get(provider)

        # Search through discovered gateway tools — match suffix after ___
        target_tool = None
        for t in gateway.gateway_tools:
            t_name = getattr(t, 'tool_name', getattr(t, 'name', str(t)))
            # Exact match or suffix match after ___
            if t_name == target_tool_name or t_name.endswith(f'___{target_tool_name}'):
                target_tool = t
                break

        if not target_tool:
            available = [getattr(t, 'tool_name', str(t)) for t in gateway.gateway_tools]
            logger.error(f"Tool for {provider} not found in gateway. Available: {available}")
            return f"Web search failed: {provider} tool not found in gateway. Available tools: {available}"

        logger.info(f"Calling gateway tool: {getattr(target_tool, 'tool_name', str(target_tool))} for {provider}")

        # Build args based on provider — MCP tool parameters
        if provider == 'tavily':
            # Tavily MCP tool params: query, max_results, search_depth, topic
            tool_args = {"query": query, "max_results": k}
        elif provider == 'serpapi':
            # SerpAPI MCP tool params: params (dict with q, engine, etc.), mode
            tool_args = {"params": {"q": query, "engine": "google_light", "num": str(k)}}
        else:
            tool_args = {"query": query}

        # Call the tool via MCP client
        # MCP client expects: call_tool_sync(tool_use_id, name, arguments)
        import uuid
        tool_use_id = f"search_{uuid.uuid4().hex[:8]}"
        result = gateway.mcp_client.call_tool_sync(
            tool_use_id,
            getattr(target_tool, 'tool_name', str(target_tool)),
            tool_args
        )

        # Format the gateway response
        if result:
            # Debug: log the result structure
            logger.info(f"MCP result type: {type(result)}")

            # Handle dict result (from strands MCP client)
            if isinstance(result, dict):
                logger.info(f"MCP result keys: {list(result.keys())}")
                logger.info(f"MCP result sample: {str(result)[:500]}")

                # Check for error status
                if result.get('status') == 'error':
                    error_msg = result.get('content', [{}])[0].get('text', str(result))
                    logger.error(f"MCP tool returned error: {error_msg}")
                    raise Exception(f"MCP tool error: {error_msg}")

                # Extract content from dict result
                content = result.get('content', [])
                if content and isinstance(content, list) and len(content) > 0:
                    if isinstance(content[0], dict):
                        result_text = content[0].get('text', str(result))
                    else:
                        result_text = str(content[0])
                else:
                    # Fallback: just stringify the whole result
                    result_text = str(result)

            # Handle object result (MCPToolResult)
            elif hasattr(result, 'content'):
                if hasattr(result, 'status') and result.status == 'error':
                    error_msg = result.content[0].get('text', str(result)) if result.content else str(result)
                    logger.error(f"MCP tool returned error: {error_msg}")
                    raise Exception(f"MCP tool error: {error_msg}")

                result_text = result.content[0].get('text', str(result)) if isinstance(result.content[0], dict) else str(result.content[0])
            else:
                result_text = str(result)

            logger.info(f"Extracted result text length: {len(result_text)}")
            if len(result_text) > 5000:
                result_text = result_text[:5000] + "\n... [truncated]"
            return f"Web Search Results (engine: {provider.title()} via Gateway)\nQuery: '{query}'\n{'='*80}\n\n{result_text}"
        else:
            logger.warning("No result returned from MCP gateway")
            return "No results returned from gateway."

    except Exception as e:
        logger.error(f"Gateway search failed for {provider}: {str(e)}")
        return f"Web search failed for {provider}: {str(e)}"


def _format_playwright_results(result: dict, query: str, include_images: bool) -> str:
    """Format Playwright Scraper search results."""
    try:
        if include_images:
            results = result.get("results", [])
            if not results:
                return "No images found."
            formatted = [f"Found {len(results)} images using {result.get('engine', 'unknown')} "
                         f"(completed in {result.get('total_time', 0):.1f}s):\n"]
            for i, item in enumerate(results, 1):
                if item.get("success"):
                    formatted.append(f"{i}. {item['url']}")
            return "\n".join(formatted)
        else:
            results = result.get("results", [])
            if not results:
                return "No results found."
            engine_used = result.get("engine", "unknown")
            total_time = result.get("total_time", 0)
            formatted = [f"Web Search Results (engine: {engine_used}, time: {total_time:.1f}s)\n",
                         f"Query: '{query}'\n", f"{'='*80}\n"]
            for i, item in enumerate(results, 1):
                if item.get("success"):
                    formatted.append(f"\n[{i}] {item.get('title', 'No title')}")
                    formatted.append(f"URL: {item.get('url', 'N/A')}")
                    formatted.append(f"Words: {item.get('word_count', 0)}")
                    content = item.get('content', '')
                    if len(content) > 1000:
                        formatted.append(f"Content: {content[:1000]}... [truncated]")
                    else:
                        formatted.append(f"Content: {content}")
                    formatted.append("-" * 80)
            return "\n".join(formatted)
    except Exception as e:
        logger.error(f"Error formatting web search results: {str(e)}")
        return f"Error formatting results: {str(e)}"


@tool
def web_search(query: str, k: int = 10, include_images: bool = False) -> str:
    """
    Search the web using the configured provider (Playwright Scraper, Tavily, or SerpAPI).
    
    - Playwright Scraper: Direct HTTP (default, always available, no config needed)
    - Tavily/SerpAPI: Routed through AgentCore Gateway (MCP targets, gateway handles auth)

    Args:
        query: The search query (1-500 characters)
        k: Number of results to return (1-20, default: 10)
        include_images: Set to True to search for images (Playwright Scraper only)

    Returns:
        Formatted search results with titles, URLs, and content
    """
    if not query or len(query) > 500:
        return "Error: Query must be between 1-500 characters"
    if k < 1 or k > 20:
        return "Error: k must be between 1 and 20"

    provider = WEB_SEARCH_PROVIDER
    logger.info(f"web_search called: provider={provider}, query='{query}', k={k}")

    if provider in ('tavily', 'serpapi'):
        return _search_via_gateway(query, k, provider)
    else:
        return _search_playwright(query, k, include_images)


# =============================================================================
# RAG TOOLS
# =============================================================================
@tool
def search_documents(query: str, top_k: int = 5) -> str:
    """
    Search uploaded documents for relevant information using RAG.
    
    Args:
        query: The search query to find relevant document content
        top_k: Number of top results to return (default: 5)
    
    Returns:
        Relevant document excerpts matching the query
    """
    try:
        rag = get_rag_instance()
        context = rag.get_context_for_query(query, top_k)
        return context
    except Exception as e:
        logger.error(f"Error in search_documents: {str(e)}")
        return f"Error searching documents: {str(e)}"


@tool
def list_uploaded_documents() -> str:
    """
    List all documents uploaded to the project's docs folder.
    
    Returns:
        List of document names and metadata
    """
    try:
        rag = get_rag_instance()
        documents = rag.list_documents()
        
        if not documents:
            return "No documents found in the project folder."
        
        doc_list = []
        for doc in documents:
            doc_list.append(f"- {doc['name']} (Size: {doc['size']} bytes)")
        
        return f"Found {len(documents)} documents:\n" + "\n".join(doc_list)
    except Exception as e:
        logger.error(f"Error listing documents: {str(e)}")
        return f"Error listing documents: {str(e)}"


@tool
def index_documents() -> str:
    """
    Index all uploaded documents for RAG search.
    This creates vector embeddings for document chunks.
    
    Returns:
        Summary of indexed documents
    """
    try:
        rag = get_rag_instance()
        results = rag.index_all_documents()
        
        if not results:
            return "No documents to index."
        
        summary = []
        total_chunks = 0
        for doc_name, chunks in results.items():
            summary.append(f"- {doc_name}: {chunks} chunks indexed")
            total_chunks += chunks
        
        return f"Indexed {len(results)} documents ({total_chunks} total chunks):\n" + "\n".join(summary)
    except Exception as e:
        logger.error(f"Error indexing documents: {str(e)}")
        return f"Error indexing documents: {str(e)}"


# =============================================================================
# AGENT CONFIG LOADER
# =============================================================================
class AgentConfigLoader:
    """Loads agent configurations from DynamoDB"""
    
    @staticmethod
    def get_config(user_id: str, project_id: str) -> Dict[str, Any]:
        """
        Retrieve agent configuration from DynamoDB
        
        Args:
            user_id: User identifier
            project_id: Project identifier
            
        Returns:
            Agent configuration dictionary
        """
        try:
            response = table.get_item(
                Key={
                    'user_id': user_id,
                    'project_id': project_id
                }
            )
            
            if 'Item' not in response:
                raise ValueError(
                    f"No configuration found for user_id={user_id}, project_id={project_id}"
                )
            
            item = response['Item']
            
            if 'response' in item and isinstance(item['response'], dict):
                config = item['response']
                logger.info(f"Successfully loaded config for user_id={user_id}, project_id={project_id}")
                logger.info(f"Config keys: {list(config.keys())}")
                
                if 'agents_config' not in config:
                    raise ValueError(f"No agents_config found in configuration")
                
                # Inject deployment_outputs (runtime_arn, gateway_url, etc.) into config
                # These are stored at the top level of the DynamoDB item, not inside 'response'
                if 'deployment_outputs' in item and isinstance(item['deployment_outputs'], dict):
                    config['deployment_outputs'] = item['deployment_outputs']
                    logger.info(f"Deployment outputs loaded: {list(item['deployment_outputs'].keys())}")

                logger.info(f"Found {len(config.get('agents_config', {}))} agents in config")
                return config
            else:
                raise ValueError(f"No 'response' field found in DynamoDB item")
            
        except ClientError as e:
            logger.error(f"DynamoDB error: {e.response['Error']['Message']}")
            raise
        except Exception as e:
            logger.error(f"Error loading config: {str(e)}")
            raise


# =============================================================================
# DYNAMIC AGENT ORCHESTRATOR (FIXED FOR MEMORY)
# =============================================================================
class DynamicAgentOrchestrator:
    """
    Dynamic orchestrator with FIXED memory implementation.
    
    CRITICAL CHANGES:
    1. Session ID and Actor ID passed per request (not at initialization)
    2. Session manager created fresh for each request with proper IDs
    3. Agents share the session manager for memory continuity
    """
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize orchestrator with agent configurations (NO session/actor here)
        
        Args:
            config: Agent configuration from DynamoDB
        """
        self.config = config
        self.agents_config = config.get('agents_config', {})
        self.coordinator_key = self._find_coordinator_key()

        # Initialize memory manager (singleton)
        self.memory_manager = get_memory_manager()

        # Initialize platform agents
        self.platform_agents_config = config.get('platform_agents', {})
        infra_config = self.platform_agents_config.get('infrastructure_agent', {})
        if infra_config.get('enabled', False):
            set_infra_config_reference(self.agents_config, GATEWAY_URL, AWS_REGION)
            logger.info("[PLATFORM] Infrastructure Agent enabled")
        else:
            logger.info("[PLATFORM] Infrastructure Agent disabled")

        security_config = self.platform_agents_config.get('security_agent', {})
        if security_config.get('enabled', False):
            set_security_config_reference(self.agents_config, GATEWAY_URL, AWS_REGION)
            logger.info("[PLATFORM] Security Agent enabled")
        else:
            logger.info("[PLATFORM] Security Agent disabled")

        observability_config = self.platform_agents_config.get('observability_agent', {})
        if observability_config.get('enabled', False):
            deployment_outputs = config.get('deployment_outputs', {})
            set_observability_config_reference(
                self.agents_config, GATEWAY_URL, AWS_REGION,
                deployment_outputs=deployment_outputs,
            )
            logger.info("[PLATFORM] Observability Agent enabled")
        else:
            logger.info("[PLATFORM] Observability Agent disabled")

        logger.info(f"Orchestrator initialized with {len(self.agents_config)} agent configs")
        logger.info(f"Memory configured: {self.memory_manager.memory_id is not None}")
    
    def _find_coordinator_key(self) -> str:
        """Find the coordinator agent key from config"""
        if COORDINATOR_KEY in self.agents_config:
            logger.info(f"Using coordinator: {COORDINATOR_KEY}")
            return COORDINATOR_KEY
        
        coordinator_patterns = [
            'coordinator', 'central_coordinator', 'orchestrator', 
            'master', 'main', 'primary', 'router'
        ]
        
        for pattern in coordinator_patterns:
            for key in self.agents_config.keys():
                if pattern in key.lower():
                    logger.info(f"Found coordinator by pattern: {key}")
                    return key
        
        if self.agents_config:
            first_key = list(self.agents_config.keys())[0]
            logger.warning(f"No coordinator found, using first agent: {first_key}")
            return first_key
        
        raise ValueError("No agents found in configuration")
    
    def _get_tools(self, tool_names: List[str], actor_id: str = '') -> List:
        """Map tool names to actual tool implementations.

        Each agent only gets the tools listed in its spec config.
        Tool name formats:
          - "web_search": Web search (playwright_scraper, tavily, or serpapi via gateway)
          - "connector:<name>": Identity-based connector (per-user OAuth, direct MCP)
          - "mcp_<service>": Load tools from a specific gateway target (e.g. "mcp_slack", "mcp_gmail")
          - "mcp": Load ALL gateway tools (backward compat)
          - "mcp:<target_name>": Load by target name (backward compat)
          - "rag_mcp_server:<target_name>": Load RAG MCP tools from gateway target
          - Other names: matched against local tool_map

        Returns:
            List of tool functions. For connectors, also stores MCP clients
            in self._connector_mcp_clients for cleanup.
        """
        tools = []

        # Import image tools
        try:
            from image_gen_tool import (
                generate_image,
                analyze_image,
                list_generated_images,
                delete_image,
                generate_image_from_image
            )
            image_tools_available = True
        except ImportError:
            logger.warning("Image generation tools not available")
            image_tools_available = False

        # Local tool map (non-gateway tools)
        tool_map = {
            'search_documents': search_documents,
            'list_uploaded_documents': list_uploaded_documents,
            'index_documents': index_documents,
            'web_search': web_search,  # Routes to playwright/tavily/serpapi
            'knowledge_base_query': knowledge_base_query,
            'knowledge_base_list_sources': knowledge_base_list_sources,
            'crawl_and_index_website': crawl_and_index_website,
            # Connector indexing tools (custom data source - S3 Vectors)
            'index_confluence_pages': index_confluence_pages,
            'index_sharepoint_documents': index_sharepoint_documents,
            'index_salesforce_records': index_salesforce_records,
            # Report generation
            'generate_report': generate_report,
        }

        if image_tools_available:
            tool_map.update({
                'generate_image': generate_image,
                'analyze_image': analyze_image,
                'list_generated_images': list_generated_images,
                'delete_image': delete_image,
                'generate_image_from_image': generate_image_from_image,
            })

        # Track which MCP targets to load
        mcp_load_all = False
        mcp_target_names = []
        connector_names = []

        # Add requested tools
        for tool_name in tool_names:
            # Connector: "connector:<name>" → Identity-based per-user MCP
            if tool_name.startswith('connector:'):
                connector_name = tool_name[10:]
                if connector_name:
                    connector_names.append(connector_name)
                continue
            if tool_name in ('mcp', 'MCP'):
                mcp_load_all = True
                continue
            # Primary format: mcp_<service> (e.g. mcp_slack, mcp_gmail)
            # The target name in the gateway matches the service name
            if tool_name.startswith('mcp_') and tool_name not in tool_map:
                target_name = tool_name[4:]  # strip "mcp_" prefix
                if target_name:
                    mcp_target_names.append(target_name)
                continue
            # Backward compat: mcp:<target_name>, mcp_server:<target_name>, integration:<target_name>
            if tool_name.startswith('mcp:'):
                target_name = tool_name[4:]
                if target_name:
                    mcp_target_names.append(target_name)
                continue
            if tool_name.startswith('integration:'):
                target_name = tool_name[12:]
                if target_name:
                    mcp_target_names.append(target_name)
                continue
            if tool_name.startswith('mcp_server:'):
                target_name = tool_name[11:]
                if target_name:
                    mcp_target_names.append(target_name)
                continue
            # rag_mcp_server:<target_name> → load tools from gateway target
            if tool_name.startswith('rag_mcp_server:'):
                target_name = tool_name[15:]
                if target_name:
                    mcp_target_names.append(target_name)
                continue
            # "rag" / "RAG" is an alias for "knowledge_base_query"
            resolved_name = 'knowledge_base_query' if tool_name in ('rag', 'RAG') else tool_name
            if resolved_name in tool_map:
                tools.append(tool_map[resolved_name])
            else:
                logger.warning(f"Tool '{tool_name}' not found in local tool map")

        # Load gateway tools based on what was requested
        if mcp_load_all or mcp_target_names:
            gateway = get_mcp_gateway()
            if gateway and gateway.is_connected:
                if mcp_load_all:
                    tools.extend(gateway.gateway_tools)
                    logger.info(f"Added ALL {len(gateway.gateway_tools)} MCP gateway tool(s)")
                else:
                    for target_name in mcp_target_names:
                        target_tools = gateway.get_tools_for_target(target_name)
                        tools.extend(target_tools)
                        logger.info(f"Added {len(target_tools)} tool(s) from target '{target_name}'")
            else:
                if mcp_load_all:
                    logger.warning("MCP requested but gateway not available")
                for tn in mcp_target_names:
                    logger.warning(f"MCP target '{tn}' requested but gateway not available")

        # Load connector tools (Identity-based per-user OAuth)
        if connector_names and actor_id:
            from connector_manager import get_connector_manager
            cm = get_connector_manager()
            if cm:
                connector_configs = [
                    _CONNECTOR_TARGET_MAP[name]
                    for name in connector_names
                    if name in _CONNECTOR_TARGET_MAP
                ]
                if connector_configs:
                    conn_tools, conn_clients, tools_by_connector = cm.get_connector_tools(connector_configs, actor_id)

                    # Filter out blocked tools per connector (scoped to avoid cross-connector conflicts)
                    blocked_tools_set = set()
                    for cfg in connector_configs:
                        cname = cfg.get('name', '')
                        cfg_blocked = cfg.get('blocked_tools', [])
                        if cfg_blocked and cname in tools_by_connector:
                            # Only block tools that actually belong to this connector
                            connector_tool_names = set(tools_by_connector[cname])
                            for bt in cfg_blocked:
                                if bt in connector_tool_names:
                                    blocked_tools_set.add(bt)

                    if blocked_tools_set:
                        before_count = len(conn_tools)
                        conn_tools = [
                            t for t in conn_tools
                            if getattr(t, 'tool_name', getattr(t, '__name__', '')) not in blocked_tools_set
                        ]
                        # Also filter tools_by_connector
                        for cname in tools_by_connector:
                            tools_by_connector[cname] = [
                                tn for tn in tools_by_connector[cname] if tn not in blocked_tools_set
                            ]
                        logger.info(f"Filtered {before_count - len(conn_tools)} blocked tool(s): {blocked_tools_set}")

                    tools.extend(conn_tools)
                    # Store clients for cleanup after request
                    if not hasattr(self, '_connector_mcp_clients'):
                        self._connector_mcp_clients = []
                    self._connector_mcp_clients.extend(conn_clients)
                    # Store per-connector tool names for prompt enhancement
                    if not hasattr(self, '_connector_tool_names'):
                        self._connector_tool_names = {}
                    self._connector_tool_names.update(tools_by_connector)
                    logger.info(f"Loaded {len(conn_tools)} connector tool(s) for user {actor_id}")
            else:
                logger.warning("Connector requested but WORKLOAD_IDENTITY_NAME not set")
        elif connector_names and not actor_id:
            logger.warning(f"Connector(s) {connector_names} requested but no actor_id provided")

        return tools
        
    def create_agents_for_request(
        self, 
        session_id: str, 
        actor_id: str
    ) -> Dict[str, Any]:
        """
        Create agents with memory for a specific request.
        
        CRITICAL FIX: Only the coordinator gets memory (session manager).
        Specialized agents run WITHOUT memory since they're called as tools.
        
        Args:
            session_id: Session ID from payload (persistent across conversation)
            actor_id: Actor/User ID from payload (for user isolation)
            
        Returns:
            Dictionary containing coordinator and specialized agents
        """
        # Create session manager with the request-specific IDs
        session_manager = None
        if self.memory_manager.memory_id:
            session_manager = self.memory_manager.create_session_manager(
                actor_id=actor_id,
                session_id=session_id
            )
            logger.info(f"[OK] Session manager created for session={session_id}, actor={actor_id}")
        else:
            logger.warning("Memory not configured. Agents will run without memory.")
        
        # Initialize specialized agents WITHOUT memory (they're tools)
        agents = {}
        agent_count = 0
        self._agent_connector_info = {}  # Track which agents have connector tools
        for agent_key, agent_config in self.agents_config.items():
            if agent_key == self.coordinator_key:
                continue  # Skip coordinator, create it separately

            try:
                tools = self._get_tools(agent_config.get('tools', []), actor_id=actor_id)
                # Track connector tools for this agent
                connector_tool_names = [t for t in agent_config.get('tools', []) if t.startswith('connector:')]
                if connector_tool_names and hasattr(self, '_connector_mcp_clients') and self._connector_mcp_clients:
                    agent_name_val = agent_config.get('name', agent_key)
                    self._agent_connector_info[agent_name_val] = [c.split(':', 1)[1] for c in connector_tool_names]

                # Enhance sub-agent prompt with connected services info (same pattern as POC agent.py)
                agent_prompt = agent_config.get('prompt', '')
                if connector_tool_names and hasattr(self, '_connector_tool_names') and self._connector_tool_names:
                    services = [c.split(':', 1)[1] for c in connector_tool_names]
                    service_map = {
                        'slack': 'Slack',
                        'google': 'Google Workspace (Gmail, Drive, Calendar)',
                        'microsoft': 'Microsoft Teams',
                    }
                    services_text = ', '.join(service_map.get(s, s) for s in services)
                    connector_only_tools = []
                    for svc in services:
                        connector_only_tools.extend(self._connector_tool_names.get(svc, []))
                    agent_prompt += (
                        f"\n\nYou have access to the user's connected services: {services_text}. "
                        f"Use the appropriate tools to help the user. "
                        f"For Slack, always pass username as 'Qubitz' when sending messages."
                    )
                    logger.info(f"Enhanced prompt for {agent_config.get('name', agent_key)} with connector tools: {connector_only_tools}")

                # Create scoped session manager if agent has memory in spec
                agent_session_manager = None
                agent_memory_config = agent_config.get('memory', 'none')
                if agent_memory_config and agent_memory_config != 'none' and self.memory_manager.memory_id:
                    agent_name_val = agent_config.get('name', agent_key)
                    scoped_session_id = f"{session_id}__{agent_name_val}"
                    agent_session_manager = self.memory_manager.create_session_manager(
                        actor_id=actor_id,
                        session_id=scoped_session_id
                    )
                    logger.info(f"[OK] Memory enabled for specialist '{agent_name_val}': {agent_memory_config} (session: {scoped_session_id})")

                agent = Agent(
                    name=agent_config.get('name', agent_key),
                    model=_fix_model_region(agent_config.get('model', 'eu.anthropic.claude-sonnet-4-5-20250929-v1:0')),
                    system_prompt=agent_prompt,
                    tools=tools,
                    session_manager=agent_session_manager,
                )

                agent_name = agent_config.get('name', agent_key)
                agents[agent_name] = agent
                agent_count += 1

                memory_status = agent_memory_config if agent_session_manager else 'disabled'
                logger.info(
                    f"[OK] Created agent [{agent_count}]: {agent_name} "
                    f"(key: {agent_key}, memory: {memory_status})"
                )
                
            except Exception as e:
                logger.error(f"[ERROR] Error creating agent {agent_key}: {str(e)}")
                raise
        
        # SECOND PASS: Wire agent_as_a_tool for specialized agents
        # Any agent with agent_as_a_tool in its config gets those agents as callable tools
        for agent_key, agent_config in self.agents_config.items():
            if agent_key == self.coordinator_key:
                continue
            agent_as_a_tool_names = agent_config.get('agent_as_a_tool', [])
            if not agent_as_a_tool_names:
                continue
            
            agent_name = agent_config.get('name', agent_key)
            if agent_name not in agents:
                continue
            
            # Build subset of agents that this agent can call
            sub_agents = {}
            for ref_name in agent_as_a_tool_names:
                if ref_name in agents and ref_name != agent_name:
                    sub_agents[ref_name] = agents[ref_name]
                else:
                    logger.warning(f"agent_as_a_tool: '{ref_name}' not found for agent '{agent_name}', skipping")
            
            if sub_agents:
                sub_agent_tools = self._create_dynamic_agent_tools(sub_agents)
                # Rebuild agent with the additional agent tools, preserving its session manager
                existing_tools = self._get_tools(agent_config.get('tools', []), actor_id=actor_id)
                all_agent_tools = existing_tools + sub_agent_tools

                # Preserve the session manager from the first pass
                existing_sm = getattr(agents[agent_name], 'session_manager', None)

                agents[agent_name] = Agent(
                    name=agent_name,
                    model=_fix_model_region(agent_config.get('model', 'eu.anthropic.claude-sonnet-4-5-20250929-v1:0')),
                    system_prompt=agent_config.get('prompt', ''),
                    tools=all_agent_tools,
                    session_manager=existing_sm,
                )
                logger.info(
                    f"Wired agent_as_a_tool for '{agent_name}': "
                    f"{list(sub_agents.keys())} ({len(sub_agent_tools)} tool wrappers), "
                    f"memory: {'enabled' if existing_sm else 'disabled'}"
                )
        
        # ── Platform Agents ────────────────────────────────────────────────
        # Create Infrastructure Agent if enabled in platform_agents config
        infra_cfg = self.platform_agents_config.get('infrastructure_agent', {})
        if infra_cfg.get('enabled', False):
            infra_tools = [
                get_system_status,
                get_recent_errors,
                get_trace_summary,
                check_model_availability,
                check_agent_tools,
            ]
            infra_model = infra_cfg.get('model', 'eu.anthropic.claude-sonnet-4-5-20250929-v1:0')
            infra_agent = Agent(
                name='infrastructure-agent',
                model=_fix_model_region(infra_model),
                system_prompt=INFRASTRUCTURE_AGENT_PROMPT,
                tools=infra_tools,
            )
            agents['infrastructure-agent'] = infra_agent
            logger.info("[PLATFORM] Infrastructure Agent created with 5 tools")

        # Create Security Agent if enabled in platform_agents config
        security_cfg = self.platform_agents_config.get('security_agent', {})
        if security_cfg.get('enabled', False):
            security_tools = [
                audit_iam_role,
                check_bedrock_guardrails,
                check_secrets_compliance,
                scan_text_for_pii,
            ]
            security_model = security_cfg.get('model', 'eu.anthropic.claude-sonnet-4-5-20250929-v1:0')
            security_agent = Agent(
                name='security-agent',
                model=_fix_model_region(security_model),
                system_prompt=SECURITY_AGENT_PROMPT,
                tools=security_tools,
            )
            agents['security-agent'] = security_agent
            logger.info("[PLATFORM] Security Agent created with 4 tools")

        # Create Observability Agent if enabled in platform_agents config
        obs_cfg = self.platform_agents_config.get('observability_agent', {})
        if obs_cfg.get('enabled', False):
            obs_tools = [
                get_runtime_metrics,
                get_session_details,
                get_user_metrics,
            ]
            obs_model = obs_cfg.get('model', 'eu.anthropic.claude-sonnet-4-5-20250929-v1:0')
            obs_agent = Agent(
                name='observability-agent',
                model=_fix_model_region(obs_model),
                system_prompt=OBSERVABILITY_AGENT_PROMPT,
                tools=obs_tools,
            )
            agents['observability-agent'] = obs_agent
            logger.info("[PLATFORM] Observability Agent created with 3 tools")

        # Re-create coordinator to pick up any re-wired agents (+ platform agents)
        coordinator = self._create_coordinator(agents, session_manager, session_id, actor_id)

        return {
            'coordinator': coordinator,
            'agents': agents,
            'session_id': session_id,
            'actor_id': actor_id,
            'memory_enabled': session_manager is not None
        }
    
    def _create_dynamic_agent_tools(self, agents: Dict[str, Any]) -> List:
        """Create agent-as-tool wrappers"""
        agent_tools = []
        
        for agent_name, agent_instance in agents.items():
            agent_config = None
            for key, cfg in self.agents_config.items():
                if cfg.get('name', key) == agent_name:
                    agent_config = cfg
                    break
            
            def make_agent_tool(agent_ref, name, description):
                """Factory function to create agent tool with proper closure.

                Uses a lock to prevent concurrent calls to the same agent instance.
                Strands Agent maintains internal conversation history, so parallel
                calls would corrupt tool_use/tool_result ID pairing and cause
                Bedrock ValidationException.
                """
                tool_name = name.replace('-', '_').replace(' ', '_')
                _lock = threading.Lock()

                def _fn(query: str) -> str:
                    """Dynamically generated agent tool"""
                    with _lock:
                        try:
                            result = agent_ref(query)
                            return str(result)
                        except Exception as e:
                            logger.error(f"Error in agent {name}: {str(e)}")
                            return f"Error executing {name}: {str(e)}"

                # Set name and doc BEFORE applying @tool so it captures the correct name
                _fn.__name__ = tool_name
                _fn.__doc__ = f"""{description}

Args:
    query: Request or query for the {name}

Returns:
    Response from the {name}
"""
                return tool(_fn)
            
            description = agent_config.get('description', f'Agent: {agent_name}') if agent_config else f'Agent: {agent_name}'
            agent_tool = make_agent_tool(agent_instance, agent_name, description)
            agent_tools.append(agent_tool)
            
            logger.info(f"[OK] Created tool: {agent_tool.__name__}")
        
        return agent_tools
    
    def _create_coordinator(
        self,
        agents: Dict[str, Any],
        session_manager,
        session_id: str,
        actor_id: str
    ) -> Agent:
        """Create coordinator agent with all tools including MCP gateway tools"""
        coordinator_config = self.agents_config.get(self.coordinator_key, {})

        # If coordinator has agent_as_a_tool, only wire those specific agents; otherwise wire all
        # Platform agents (e.g. infrastructure-agent) are ALWAYS included regardless of filter
        platform_agent_names = {'infrastructure-agent', 'security-agent', 'observability-agent'}
        coordinator_agent_refs = coordinator_config.get('agent_as_a_tool', [])
        if coordinator_agent_refs:
            filtered_agents = {
                name: inst for name, inst in agents.items()
                if name in coordinator_agent_refs or name in platform_agent_names
            }
            agent_tools = self._create_dynamic_agent_tools(filtered_agents)
            logger.info(f"Coordinator agent_as_a_tool: wiring {list(filtered_agents.keys())} (filtered from {len(agents)} agents)")
        else:
            agent_tools = self._create_dynamic_agent_tools(agents)
        additional_tools = self._get_tools(coordinator_config.get('tools', []), actor_id=actor_id)

        # Image tools — only if coordinator spec includes image tools
        image_tools = []
        coordinator_tool_names = coordinator_config.get('tools', [])
        try:
            from image_gen_tool import (
                generate_image,
                analyze_image,
                list_generated_images,
                delete_image,
                generate_image_from_image
            )
            image_tool_names = ['generate_image', 'analyze_image', 'list_generated_images', 'delete_image', 'generate_image_from_image']
            if any(t in coordinator_tool_names for t in image_tool_names):
                image_tools = [
                    generate_image,
                    analyze_image,
                    list_generated_images,
                    delete_image,
                    generate_image_from_image
                ]
        except ImportError:
            logger.warning("Image tools not available for coordinator")

        # RAG tools — only if coordinator spec includes RAG tools
        rag_tools = []
        rag_tool_names = ['search_documents', 'list_uploaded_documents', 'index_documents']
        if any(t in coordinator_tool_names for t in rag_tool_names):
            rag_tools = [search_documents, list_uploaded_documents, index_documents]

        # Knowledge Base tools — only if coordinator spec includes KB tools
        kb_tools = []
        kb_tool_names = ['knowledge_base_query', 'knowledge_base_list_sources']
        if any(t in coordinator_tool_names for t in kb_tool_names):
            kb_tools = [knowledge_base_query, knowledge_base_list_sources]

        # Web search (Playwright Scraper) — only if "web_search" is in coordinator's tools
        web_tools = []
        if 'web_search' in coordinator_tool_names:
            web_tools = [web_search]

        # Gateway tools (Tavily/SerpAPI) — only if "mcp" is in coordinator's tools
        gateway_tools = []
        gateway = None
        if 'mcp' in coordinator_tool_names:
            gateway = get_mcp_gateway()
            if gateway and gateway.is_connected:
                gateway_tools = gateway.gateway_tools

        # Combine all tools
        all_tools = agent_tools + additional_tools + rag_tools + kb_tools + web_tools + image_tools + gateway_tools
        
        coordinator_prompt = coordinator_config.get('prompt', '')
        
        if agent_tools or image_tools:
            agent_list = '\n'.join([
                f"- {agent_name}: {self.agents_config.get(key, {}).get('description', 'N/A')}"
                for key, cfg in self.agents_config.items()
                for agent_name in [cfg.get('name', key)]
                if key != self.coordinator_key
            ])
            
            memory_info = ""
            if session_manager:
                memory_status = "ENABLED"
                capabilities = []
                if ENABLE_SHORT_TERM_MEMORY:
                    capabilities.append("short-term (conversation history within session)")
                if ENABLE_LONG_TERM_MEMORY:
                    capabilities.append("long-term (user preferences & facts across sessions)")
                
                memory_capabilities = " and ".join(capabilities) if capabilities else "configured but disabled"
                
                memory_info = f"""

    MEMORY CAPABILITIES ({memory_status}):
    You have {memory_capabilities}.
    - Actor ID: {actor_id} (user isolation - remember things specific to THIS user)
    - Session ID: {session_id} (conversation continuity - maintain context within THIS conversation)

    IMPORTANT MEMORY USAGE:
    - The conversation history is automatically maintained - you have full context of this session
    - User preferences and facts are automatically retrieved from long-term memory
    - You can reference earlier parts of the conversation naturally
    - Information learned about this user persists across their different sessions
    - When you delegate to specialized agents, provide them with necessary context from the conversation
    """
            
            # Build connector info for prompt
            connector_info = ""
            if hasattr(self, '_agent_connector_info') and self._agent_connector_info:
                connector_lines = []
                for agent_name, connectors in self._agent_connector_info.items():
                    connector_lines.append(f"    - {agent_name}: connected to {', '.join(c.title() for c in connectors)}")
                connector_info = f"""

    CONNECTED SERVICES (per-user OAuth connectors):
    The following agents have access to the user's connected services:
{chr(10).join(connector_lines)}
    When the user asks about these services, delegate to the appropriate agent.
    These connectors use the user's own OAuth tokens - each user has their own access.
"""

            enhanced_prompt = f"""{coordinator_prompt}

    You have access to the following specialized agents as tools:
    {agent_list}
    {connector_info}
    You also have powerful IMAGE GENERATION and ANALYSIS tools:
    - generate_image: Create images from text descriptions using Stability AI models
    * Models available: 'sd3-large' (default, balanced), 'stable-image-ultra' (highest quality), 
        'stable-image-core' (fast & efficient)
    * Can specify negative prompts, dimensions, and random seed
    * Returns S3 URI and HTTPS URL to the generated image
    
    - analyze_image: Analyze any image using Claude's vision capabilities
    * Can analyze images from S3 URIs, S3 keys, or local paths
    * Ask questions about image content, objects, text, scenes, etc.
    * Excellent for document analysis, object detection, image understanding
    
    - generate_image_from_image: Transform existing images (image-to-image)
    * Modify existing images based on text prompts
    * Control transformation strength (0.0 = no change, 1.0 = completely new)
    * Great for style transfer, modifications, variations
    
    - list_generated_images: See all images generated in this project
    
    - delete_image: Remove unwanted generated images

    IMPORTANT: All generated images are automatically stored in S3 at:
    s3://qubitz-customer-prod-v2/{USER_ID}/{PROJECT_ID}/images/

    You also have RAG (Retrieval Augmented Generation) tools for searching uploaded documents:
    - search_documents: Search uploaded documents for relevant information using semantic search
    - list_uploaded_documents: List all documents uploaded to the project folder
    - index_documents: Index all uploaded documents for RAG search (creates vector embeddings)

    Documents are stored at: s3://{S3_DOCS_BUCKET}/{USER_ID}/{PROJECT_ID}/

    You also have KNOWLEDGE BASE tools using CUSTOM DATA SOURCE (S3 Vectors):
    IMPORTANT: All connectors use S3 Vectors for indexing, NOT Bedrock Knowledge Base.
    This provides full control over data ingestion and retrieval.
    
    - knowledge_base_query: Search across ALL indexed data sources (S3, web crawler, Confluence,
      SharePoint, Salesforce) using S3 Vectors custom data source for this user/project.
      Returns relevant content with source attribution and relevance scores.
    - knowledge_base_list_sources: List all documents indexed in the knowledge base.
    - crawl_and_index_website: Crawl a website and index pages into S3 Vectors.
    - index_confluence_pages: Fetch Confluence pages via API and index into S3 Vectors.
    - index_sharepoint_documents: Fetch SharePoint documents via Microsoft Graph API and index into S3 Vectors.
    - index_salesforce_records: Fetch Salesforce records via REST API and index into S3 Vectors.

    KNOWLEDGE BASE vs DOCUMENT SEARCH:
    - Use knowledge_base_query for data from external systems (Confluence, SharePoint, S3 buckets, websites)
    - Use search_documents for user-uploaded files (PDF, DOCX, CSV, etc.)
    - Both can be used together for comprehensive information retrieval
    - All external data is fetched via APIs and indexed into S3 Vectors (custom data source)

    You have access to WEB SEARCH capabilities:"""

            # Add MCP gateway tools info if connected
            if gateway and gateway.is_connected:
                gateway_tool_names = [getattr(t, 'tool_name', str(t)) for t in gateway_tools]
                enhanced_prompt += f"""

    MCP GATEWAY TOOLS (auto-discovered from AgentCore Gateway):
    The following tools are available through the gateway: {', '.join(gateway_tool_names)}
    These tools are provided by gateway targets (e.g., Tavily Search, SerpAPI Search) and are
    called through the AgentCore Gateway MCP protocol. Use them as regular tools - the gateway
    handles authentication and routing automatically.
"""

            enhanced_prompt += """
    - web_search: Playwright Scraper (direct, always available)
      * Uses Quick API first (fast, 4-6s) with automatic Primary API fallback (8-15s if needed)
      * Returns 1-20 results with full text content, titles, and URLs
      * Multi-engine fallback: Quick API (DuckDuckGo/Brave) -> Primary API (Brave/Startpage/Yahoo/Yandex)
      * Set include_images=True to search for images instead of web content
      * Example: web_search(query="latest AI developments", k=10)

    WEB SEARCH BEST PRACTICES:
    1. Use gateway tools (Tavily/Brave) when available for best quality results
    2. Use web_search (Playwright Scraper) as fallback or for image search
    3. Start with 5-10 results for balanced speed and coverage
    4. Use specific, clear queries for better results
    {memory_info}

    IMPORTANT: You (the coordinator) have memory, but the specialized agents you call do NOT have memory.
    When delegating to specialized agents, include relevant context from the conversation in your query to them.

    Use these agents, tools, and memory intelligently to fulfill user requests. Delegate tasks to the most appropriate agent(s).

    IMAGE GENERATION BEST PRACTICES:
    1. Be specific and descriptive in prompts for better results
    2. Use negative prompts to avoid unwanted elements
    3. For photorealistic images, use 'stable-image-ultra'
    4. For faster generation, use 'stable-image-core'
    5. Always provide users with both S3 URI and HTTPS URL
    6. If image generation fails due to content filtering, suggest alternative phrasing

    """
        else:
            enhanced_prompt = coordinator_prompt
        
        # Create coordinator with all tools
        coordinator = Agent(
            name=coordinator_config.get('name', 'coordinator'),
            model=_fix_model_region(coordinator_config.get('model', 'eu.anthropic.claude-sonnet-4-5-20250929-v1:0')),
            system_prompt=enhanced_prompt,
            tools=all_tools,
            session_manager=session_manager,
        )
        
        logger.info(
            f"Coordinator created with {len(agent_tools)} agent tools + "
            f"{len(additional_tools)} direct tools + {len(rag_tools)} RAG tools + "
            f"{len(kb_tools)} KB tools + "
            f"{len(web_tools)} web tools + {len(image_tools)} image tools + "
            f"{len(gateway_tools)} MCP gateway tools, "
            f"memory: {session_manager is not None}"
        )
        
        return coordinator
        

    def _cleanup_connector_clients(self):
        """Clean up per-request connector MCP clients."""
        if hasattr(self, '_connector_mcp_clients'):
            for mc in self._connector_mcp_clients:
                try:
                    mc.stop()
                except Exception:
                    pass
            self._connector_mcp_clients = []


# =============================================================================
# GLOBAL ORCHESTRATOR INSTANCE
# =============================================================================
orchestrator = None


def initialize_orchestrator():
    """
    Initialize the orchestrator from DynamoDB config (ONE TIME ONLY).
    
    CRITICAL: Does NOT take session_id or actor_id - those come per request!
    """
    global orchestrator
    
    try:
        if not USER_ID or not PROJECT_ID:
            raise ValueError("USER_ID and PROJECT_ID environment variables must be set")
        
        # Load config from DynamoDB
        config_loader = AgentConfigLoader()
        config = config_loader.get_config(USER_ID, PROJECT_ID)
        
        # Create orchestrator (NO session/actor IDs here!)
        orchestrator = DynamicAgentOrchestrator(config=config)
        
        logger.info("=" * 80)
        logger.info("Orchestrator initialized successfully")
        logger.info(f"Coordinator: {orchestrator.coordinator_key}")
        logger.info(f"Agent configs loaded: {len(orchestrator.agents_config)}")
        logger.info(f"Memory ID: {orchestrator.memory_manager.memory_id}")
        logger.info(f"Short-term memory: {ENABLE_SHORT_TERM_MEMORY}")
        logger.info(f"Long-term memory: {ENABLE_LONG_TERM_MEMORY}")
        logger.info("=" * 80)
        
        return orchestrator
    
    except Exception as e:
        logger.error(f"Failed to initialize orchestrator: {str(e)}")
        raise


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================
def extract_response_text(response) -> str:
    """Extract text content from agent response object"""
    if isinstance(response, str):
        return response
    
    if hasattr(response, 'message'):
        message = response.message
        if isinstance(message, dict) and 'content' in message:
            content = message['content']
            if isinstance(content, list):
                text_parts = []
                for block in content:
                    if isinstance(block, dict) and 'text' in block:
                        text_parts.append(block['text'])
                    elif isinstance(block, str):
                        text_parts.append(block)
                return '\n'.join(text_parts)
            elif isinstance(content, str):
                return content
        elif isinstance(message, str):
            return message
    
    if isinstance(response, dict):
        if 'content' in response:
            content = response['content']
            if isinstance(content, list):
                text_parts = []
                for block in content:
                    if isinstance(block, dict) and 'text' in block:
                        text_parts.append(block['text'])
                    elif isinstance(block, str):
                        text_parts.append(block)
                return '\n'.join(text_parts)
            elif isinstance(content, str):
                return content
        if 'text' in response:
            return response['text']
    
    return str(response)


def serialize_agent_info(agent_info: Dict[str, Any]) -> Dict[str, Any]:
    """Serialize agent info to JSON-compatible format"""
    return {
        'session_id': agent_info.get('session_id', ''),
        'actor_id': agent_info.get('actor_id', ''),
        'memory_enabled': agent_info.get('memory_enabled', False),
        'agent_count': len(agent_info.get('agents', {})),
        'agent_names': list(agent_info.get('agents', {}).keys())
    }


# =============================================================================
# AGENTCORE ENTRYPOINT (FIXED)
# =============================================================================
@app.entrypoint
async def invoke(payload: Dict[str, Any], context=None):
    """AgentCore entrypoint — normal JSON response or SSE streaming (stream=true in payload)."""
    s3log = get_s3_logger()
    stream_mode = payload.get('stream', False)
    
    try:
        # Extract parameters from payload
        user_prompt = payload.get('prompt', '')
        if not user_prompt:
            return {"error": "No prompt provided in payload", "status": "failed"}
        
        # CRITICAL: Get session_id and actor_id from payload (for proper isolation)
        # Option 1: From payload (recommended)
        session_id = payload.get('session_id')
        actor_id = payload.get('actor_id')
        
        # Option 2: From context if using AgentCore Runtime
        if context:
            if not session_id:
                session_id = getattr(context, 'session_id', None)
            if not actor_id:
                # Extract from custom headers if set
                headers = getattr(context, 'headers', {})
                actor_id = headers.get('x-amzn-bedrock-agentcore-runtime-custom-actorid')
        
        # Option 3: Generate defaults ONLY for testing (NOT for production)
        if not session_id:
            session_id = f"session_{datetime.utcnow().strftime('%Y%m%d_%H%M%S_%f')}"
            logger.warning(f"[WARN] No session_id in payload, generated: {session_id}")
        
        if not actor_id:
            actor_id = USER_ID
            logger.warning(f"[WARN] No actor_id in payload, using default: {actor_id}")
        
        s3log.info("=" * 80)
        s3log.info(f"Request received - Session: {session_id}, Actor: {actor_id}")
        s3log.info(f"Prompt: {user_prompt[:200]}...")
        s3log.info("=" * 80)
        
        # Initialize orchestrator (if not already done)
        global orchestrator
        if orchestrator is None:
            orchestrator = initialize_orchestrator()
        
        # CRITICAL: Create agents with proper session_id and actor_id for THIS request
        request_context = orchestrator.create_agents_for_request(
            session_id=session_id,
            actor_id=actor_id
        )
        
        coordinator = request_context['coordinator']
        
        s3log.info(f"Agents created: {len(request_context['agents'])}, Memory: {request_context['memory_enabled']}")

        if stream_mode:
            # Streaming — return generator; framework auto-sends as SSE
            return _stream_response(
                coordinator, user_prompt, session_id, actor_id,
                request_context, orchestrator, s3log
            )

        # Normal — single JSON response
        try:
            response = coordinator(user_prompt)
        finally:
            # Clean up per-request connector MCP clients (stop matches start)
            if orchestrator and hasattr(orchestrator, '_connector_mcp_clients'):
                for mc in orchestrator._connector_mcp_clients:
                    try:
                        mc.stop()
                    except Exception:
                        pass
                orchestrator._connector_mcp_clients = []

        # Extract text content from response
        response_text = extract_response_text(response)

        s3log.info(f"Response generated (length: {len(response_text)} chars)")
        s3log.info("=" * 80)
        s3log.flush()

        return {
            "result": response_text,
            "status": "success",
            "metadata": {
                "session_id": session_id,
                "actor_id": actor_id,
                "memory_enabled": request_context['memory_enabled'],
                "short_term_memory": ENABLE_SHORT_TERM_MEMORY,
                "long_term_memory": ENABLE_LONG_TERM_MEMORY,
                "agents_created": len(request_context['agents'])
            }
        }

    except Exception as e:
        s3log.error(f"[ERROR] Error processing request: {str(e)}")
        s3log.error("=" * 80)
        s3log.flush()
        logger.error(f"Error processing request: {str(e)}", exc_info=True)
        return {
            "error": str(e),
            "status": "failed"
        }


# =============================================================================
# STREAMING RESPONSE GENERATOR
# =============================================================================
async def _stream_response(coordinator, user_prompt, session_id, actor_id, request_context, orch, s3log):
    """Async generator that yields dicts — framework converts each to SSE `data: {...}`.

    Uses Agent.stream_async() to consume Strands events as they arrive.
    """
    yield {"type": "stream_start", "session_id": session_id}

    full_response = ""
    try:
        async for event in coordinator.stream_async(user_prompt):
            if "data" in event:
                chunk = event["data"]
                full_response += chunk
                yield {"type": "stream_chunk", "data": chunk}
            elif "current_tool_use" in event and isinstance(event.get("current_tool_use"), dict) and event["current_tool_use"].get("name"):
                yield {"type": "tool_use", "tool": event["current_tool_use"]["name"]}
    except Exception as e:
        logger.error(f"Streaming error: {e}", exc_info=True)
        yield {"type": "error", "message": str(e)}
        return
    finally:
        if orch and hasattr(orch, '_connector_mcp_clients'):
            for mc in orch._connector_mcp_clients:
                try:
                    mc.stop()
                except Exception:
                    pass
            orch._connector_mcp_clients = []

    s3log.info(f"Streamed response ({len(full_response)} chars)")
    s3log.flush()

    yield {
        "type": "stream_complete",
        "result": full_response,
        "metadata": {
            "session_id": session_id,
            "actor_id": actor_id,
            "memory_enabled": request_context['memory_enabled'],
            "agents_created": len(request_context['agents'])
        }
    }


# =============================================================================
# HEALTH CHECK
# =============================================================================
@app.ping
def ping() -> PingStatus:
    """Health check endpoint for AgentCore"""
    return PingStatus.HEALTHY


if __name__ == "__main__":
    import uvicorn

    host = os.getenv("HOST", "0.0.0.0")
    port = int(os.getenv("PORT", 8080))
    
    logger.info(f"ðŸš€ Multi-Agent System with AgentCore Memory")
    logger.info(f"   User: {USER_ID}")
    logger.info(f"   Project: {PROJECT_ID}")
    logger.info(f"   Region: {AWS_REGION}")
    logger.info(f"   DynamoDB: {DYNAMODB_TABLE_NAME}")
    logger.info(f"   S3 Docs Bucket: {S3_DOCS_BUCKET}")
    logger.info(f"   S3 Vector Bucket: {S3_VECTOR_BUCKET}")
    logger.info(f"   Docs Path: {S3_DOCS_BUCKET}/{USER_ID}/{PROJECT_ID}/")
    logger.info(f"   Logs Path: {S3_DOCS_BUCKET}/{USER_ID}/{PROJECT_ID}/logs.txt")
    logger.info(f"   Memory ID: {MEMORY_ID if MEMORY_ID else 'Not configured'}")
    logger.info(f"   Memory Region: {MEMORY_REGION}")
    logger.info(f"   Short-term Memory: {'Enabled' if ENABLE_SHORT_TERM_MEMORY else 'Disabled'}")
    logger.info(f"   Long-term Memory: {'Enabled' if ENABLE_LONG_TERM_MEMORY else 'Disabled'}")
    logger.info(f"   Gateway URL: {GATEWAY_URL if GATEWAY_URL else 'Not configured'}")
    logger.info(f"   Knowledge Base ID: {KNOWLEDGE_BASE_ID if KNOWLEDGE_BASE_ID else 'Not configured'}")

    uvicorn.run(app, host=host, port=port, log_level="info")


# =============================================================================
# USAGE DOCUMENTATION
# =============================================================================
"""
CLIENT-SIDE IMPLEMENTATION EXAMPLE:
===================================

import uuid
import requests

# Start new conversation - generate session_id ONCE
session_id = str(uuid.uuid4())
actor_id = "user_12345"  # From authentication (JWT, OAuth, etc.)

agent_url = "https://your-agent-endpoint.com/invoke"

# First message in conversation
response1 = requests.post(agent_url, json={
    "prompt": "Hello, my name is Alice and I like pizza",
    "session_id": session_id,
    "actor_id": actor_id
})
print(response1.json()['result'])

# Second message - SAME session_id for memory continuity
response2 = requests.post(agent_url, json={
    "prompt": "What is my name and what do I like?",
    "session_id": session_id,  # <- SAME session_id
    "actor_id": actor_id
})
print(response2.json()['result'])
# Should remember: "Your name is Alice and you like pizza"

# Third message - still same conversation
response3 = requests.post(agent_url, json={
    "prompt": "Can you recommend a restaurant?",
    "session_id": session_id,  # <- SAME session_id
    "actor_id": actor_id
})
print(response3.json()['result'])
# Should recommend pizza restaurants based on preference


# NEW CONVERSATION - different session_id
new_session_id = str(uuid.uuid4())

response4 = requests.post(agent_url, json={
    "prompt": "What do I like to eat?",
    "session_id": new_session_id,  # <- NEW session_id
    "actor_id": actor_id  # <- SAME actor_id
})
# Long-term memory should still remember: "You like pizza"


MEMORY RESOURCE CREATION (ONE-TIME SETUP):
==========================================

from bedrock_agentcore.memory import MemoryClient

client = MemoryClient(region_name="eu-central-1")

# Create memory resource with all strategies
memory = client.create_memory_and_wait(
    name="MyAgentMemory",
    description="Memory for multi-agent system with STM and LTM",
    strategies=[
        {
            "summaryMemoryStrategy": {
                "name": "SessionSummarizer",
                "namespaces": ["/summaries/{actorId}/{sessionId}"]
            }
        },
        {
            "userPreferenceMemoryStrategy": {
                "name": "PreferenceLearner",
                "namespaces": ["/preferences/{actorId}"]
            }
        },
        {
            "semanticMemoryStrategy": {
                "name": "FactExtractor",
                "namespaces": ["/facts/{actorId}"]
            }
        }
    ]
)

# Save memory_id to environment variable
memory_id = memory.get('id')
print(f"Memory ID: {memory_id}")
print(f"Set environment variable: AGENTCORE_MEMORY_ID={memory_id}")


TESTING MEMORY:
===============

# Test 1: Short-term memory (within session)
session_id = "test-session-001"
actor_id = "test-user-001"

invoke({"prompt": "I like red color", "session_id": session_id, "actor_id": actor_id})
invoke({"prompt": "What color do I like?", "session_id": session_id, "actor_id": actor_id})
# Should remember: red

# Test 2: Long-term memory (across sessions)
new_session = "test-session-002"

invoke({"prompt": "What color do I like?", "session_id": new_session, "actor_id": actor_id})
# Should still remember: red (from LTM)

# Test 3: User isolation (different actor)
different_user = "test-user-002"

invoke({"prompt": "What color do I like?", "session_id": session_id, "actor_id": different_user})
# Should NOT know - different user


KEY POINTS:
===========

1. [OK] session_id must persist across requests in the same conversation
2. [OK] actor_id identifies the user (for isolation)
3. [OK] Generate session_id once per conversation on client-side
4. [OK] Get actor_id from authentication (JWT, OAuth, etc.)
5. [OK] Memory automatically maintains conversation history
6. [OK] Long-term memory persists across sessions for the same actor
7. [OK] Different actors have completely isolated memory
"""
# deploy_timestamp: 2026-03-06T15:12:09.812510
